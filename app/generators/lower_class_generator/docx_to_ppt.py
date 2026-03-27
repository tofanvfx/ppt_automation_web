import re
import os
import copy
import io
from docx import Document
from docx.document import Document as _Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml import parse_xml
from PIL import Image


def generate_ppt(docx_path, template_path, output_path, progress_callback=None):
    def report_progress(percentage, status):
        if progress_callback:
            progress_callback(percentage, status)

    report_progress(2, "Initializing...")
    prs = Presentation(template_path)
    doc = Document(docx_path)

    def clean(text):
        return re.sub(r'[^a-z0-9]', '', text.strip().lower())

    global_metadata = {
        'class': '',
        'subject': '',
        'chapter number': '',
        'chapter name': '',
        'lesson': '',
        'topic': '',
        'subtopic': ''
    }

    def get_text(entry):
        """Extract plain text from a content entry (either a string or a (text_parts, ilvl) tuple)."""
        data = entry[0] if isinstance(entry, tuple) else entry
        if isinstance(data, list):
            # Concatenate only text parts for plain text version
            return "".join(p['value'] if p['type'] == 'text' else '' for p in data)
        return str(data)

    def copy_image_rels(element, source_part, target_part):
        """Copy image relationships referenced by blip elements from source to target part."""
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        r_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
        for blip in element.findall('.//a:blip', namespaces=ns):
            old_rid = blip.get(f'{r_ns}embed')
            if old_rid and old_rid in source_part.rels:
                rel = source_part.rels[old_rid]
                # Add the image to the target part and get a new rId
                new_rid = target_part.relate_to(rel.target_part, rel.reltype)
                blip.set(f'{r_ns}embed', new_rid)

    def remove_locks(element):
        """Remove a:spLocks and a:grpSpLocks elements from a shape XML to make it fully editable."""
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        for lock in element.findall('.//a:spLocks', namespaces=ns):
            lock.getparent().remove(lock)
        for lock in element.findall('.//a:grpSpLocks', namespaces=ns):
            lock.getparent().remove(lock)
        return element

    def apply_locks(element):
        """Add protection locks to all shapes within an element to prevent interaction and show lock icon in Selection Pane."""
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
              'p': 'http://schemas.openxmlformats.org/officeDocument/2006/main',
              'p_main': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        
        # We need to find all types of non-visual properties: shapes, pictures, groups, connectors
        # Presentation elements (p:pic, p:sp, etc.) use the presentationml namespace
        targets = [
            ('.//p_main:nvSpPr', 'p_main:cNvSpPr', 'a:spLocks'),
            ('.//p_main:nvPicPr', 'p_main:cNvPicPr', 'a:picLocks'),
            ('.//p_main:nvGrpSpPr', 'p_main:cNvGrpSpPr', 'a:grpSpLocks'),
            ('.//p_main:nvCxnSpPr', 'p_main:cNvCxnSpPr', 'a:cxnSpLocks')
        ]
        
        for search_xpath, cNvPr_tag, lock_tag in targets:
            # Find all candidates (root + descendants)
            candidates = []
            # Check descendants
            candidates.extend(element.findall(search_xpath, namespaces=ns))
            # Check root
            root_tag_local = search_xpath.split(':')[-1]
            if f"{{{ns['p_main']}}}{root_tag_local}" == element.tag:
                candidates.append(element)

            for nvPr in candidates:
                cNvPr = nvPr.find(cNvPr_tag, namespaces=ns)
                if cNvPr is not None:
                    locks = cNvPr.find(lock_tag, namespaces=ns)
                    if locks is None:
                        # Create new locks element
                        locks = parse_xml(f'<a:{lock_tag.split(":")[1]} xmlns:a="{ns["a"]}" noGrp="1" noSelect="1" noRot="1" noChangeAspect="1" noMove="1" noResize="1" noEditPoints="1" noAdjustHandles="1" noChangeArrowheads="1" noChangeShapeType="1" noTextEdit="1"/>')
                        cNvPr.append(locks)
                    else:
                        # Update existing locks
                        for attr in ("noGrp", "noSelect", "noRot", "noChangeAspect", "noMove", "noResize", "noEditPoints", "noAdjustHandles", "noChangeArrowheads", "noChangeShapeType", "noTextEdit"):
                            locks.set(attr, "1")
        return element

    def get_layout(name):
        for layout in prs.slide_layouts:
            if layout.name == f"LAYOUT_{name}":
                return layout
            
        # Try case-insensitive, underscore-insensitive match
        target = f"layout_{name}".replace("_", "").lower()
        for layout in prs.slide_layouts:
            if layout.name.replace("_", "").replace(" ", "").lower() == target:
                return layout
            
        if name.replace("_", "").lower() == 'mathtitlepage':
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_math_page_title':
                    return layout

        if name.replace("_", "").lower() == 'ssttitlepage':
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_page_title':
                    return layout
                
        if name.replace("_", "").lower() in ('learningobjective', 'sstlopage', '1layoutsstlopage'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_lo_page' or layout.name == 'LAYOUT_sst_lo_page':
                    return layout

        if name.replace("_", "").lower() == 'mathlopage':
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_math_lo_page':
                    return layout

        if name.replace("_", "").lower() in ('finalquizpage', 'finalquiz', 'quiz'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_final_quiz_page':
                    return layout

        if name.replace("_", "").lower() in ('quiztimepage', 'sstquiztimepage', 'sst_quiztime_page'):
            # Return _01 as default; the main loop will override based on option lengths
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_quiztime_page_01':
                    return layout

        if name.replace("_", "").lower() in ('sstcontentpage01', 'sstcontentpage1'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_content_page_01' or layout.name == '1_LAYOUT_sst_content_page_01':
                    return layout

        if name.replace("_", "").lower() in ('sstsummarypage', 'sstsummary', 'sstsummerypage', 'sstsummery'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_summary_page':
                    return layout

        if name.replace("_", "").lower() in ('mathsummarypage', 'mathsummary', 'mathsummerypage', 'mathsummery'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_math_summary_page':
                    return layout

        if name.replace("_", "").lower() in ('notedownpage', 'sstnotedownpage', 'sst_notedown_page'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_notedown_page':
                    return layout

        if name.replace("_", "").lower() in ('previouspage', 'sstpreviouspage', 'sst_previous_page'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_previous_page':
                    return layout

        if name.replace("_", "").lower() in ('homeworkpage', 'ssthomeworkpage', 'ssthomework', 'homework'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_homework_page':
                    return layout

        if name.replace("_", "").lower() in ('discussionpage', 'sstdiscussionpage', 'sstdiscussion', 'sst_discussion_page'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_discussion_page':
                    return layout

        if name.replace("_", "").lower() in ('mathdefaultpage', 'math_default_page'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_math_default_page':
                    return layout

        if name.replace("_", "").lower() in ('sstdeafultpage', 'sst_deafult_page'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_deafult_page':
                    return layout

        if name.replace("_", "").lower() in ('homeworkquestionpage', 'ssthomeworkquestionpage'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_homework_question_page':
                    return layout

        if name.replace("_", "").lower() in ('sstactivitystaticpage', 'sst_activity_static_page', 'activitystaticpage'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_activity_static_page':
                    return layout

        return None

    # Nested dictionary to store templates per layout
    sst_content_templates = {
        'LAYOUT_sst_content_page_01': {'topic': None, 'subtopic': None, 'text': [], 'static_elements': []},
        '1_LAYOUT_sst_content_page_01': {'topic': None, 'subtopic': None, 'text': [], 'static_elements': []},
        'LAYOUT_sst_content_page_02': {'topic': None, 'subtopic': None, 'text': [], 'static_elements': []},
        'LAYOUT_sst_notedown_page': {'topic': None, 'subtopic': None, 'text': [], 'static_elements': []},
        'LAYOUT_sst_quiztime_page_01': {'title': None, 'question': None, 'options': [], 'picture': None},
        'LAYOUT_sst_quiztime_page_02': {'title': None, 'question': None, 'options': [], 'picture': None},
        'LAYOUT_sst_discussion_page': {'question1': None, 'static_elements': []},
        'LAYOUT_sst_homework_page': {'static_elements': []},
        'LAYOUT_homework_question_page': {'text': [], 'static_elements': []},
        'LAYOUT_syr': {'static_elements': []},
        'LAYOUT_ask_question': {'static_elements': []},
        'LAYOUT_sst_activity_page_01': {'static_elements': [], 'text': []},
        'LAYOUT_sst_activity_page_02': {'static_elements': [], 'text': []},
        'LAYOUT_sst_deafult_page': {'topic': None, 'subtopic': None, 'text': [], 'static_elements': []},
        'LAYOUT_math_default_page': {'topic': None, 'subtopic': None, 'text': [], 'static_elements': []},
        'LAYOUT_sst_activity_static_page': {'static_elements': []}
    }

    # Extract lo_page group shape XML before processing slides
    lo_group_xmls = {
        'LAYOUT_sst_lo_page': None,
        'LAYOUT_math_lo_page': None,
        'LAYOUT_sst_summary_page': None,
        'LAYOUT_math_summary_page': None,
        'LAYOUT_sst_previous_page': None
    }

    # Store static title templates for lo/summary layouts where add_slide fails
    lo_title_xmls = {k: None for k in lo_group_xmls.keys()}
    lo_title_xmls['LAYOUT_sst_lo_page'] = None # Ensure renamed one is included
    lo_title_xmls['LAYOUT_sst_previous_page'] = None

    # Store subtitle paragraph XML templates per layout (list of para XMLs for ilvl 1,2,3)
    lo_subtitle_para_xmls = {k: [] for k in lo_group_xmls.keys()}

    # Capture logo elements to be added to every slide
    logo_elements = []
    logo_source_part = None

    def inject_logo(target_slide):
        """Inject captured logo elements onto a target slide and lock them."""
        for logo_xml in logo_elements:
            new_logo_elem = apply_locks(copy.deepcopy(logo_xml))
            target_slide.shapes._spTree.append(new_logo_elem)
            if logo_source_part:
                copy_image_rels(new_logo_elem, logo_source_part, target_slide.part)

    for layout in prs.slide_layouts:
        if layout.name in lo_group_xmls:
            for shape in layout.shapes:
                if shape.shape_type == 6:  # msoGroup
                    for subshape in shape.shapes:
                        if getattr(subshape, 'has_text_frame', False) and 'Text goes here' in subshape.text:
                            # Capture subtitle paragraph XML templates (paragraphs after the first)
                            tf = subshape.text_frame
                            for pi in range(1, len(tf.paragraphs)):
                                lo_subtitle_para_xmls[layout.name].append(
                                    copy.deepcopy(tf.paragraphs[pi]._p)
                                )
                            lo_group_xmls[layout.name] = copy.deepcopy(shape.element)
                            shape.element.getparent().remove(shape.element)
                            break
                    if lo_group_xmls[layout.name] is not None:
                        break
    
        if layout.name == "LAYOUT_logo":
            logo_source_part = layout.part
            for shape in list(layout.shapes):
                logo_elements.append(copy.deepcopy(shape.element))
                shape.element.getparent().remove(shape.element)
    
        if layout.name in lo_title_xmls:
            for shape in layout.shapes:
                if shape.shape_type != 6 and getattr(shape, 'has_text_frame', False) and shape.text.strip():
                    lo_title_xmls[layout.name] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                    # print(f"DEBUG: Captured title for {layout.name}: '{shape.text[:20]}...'")
                    break
    
        if layout.name in sst_content_templates:
            templates = sst_content_templates[layout.name]
            for shape in list(layout.shapes): # Use list() to avoid issues when removing
                has_topic = False
                has_subtopic = False
                has_quiz_title = False
                has_body_text = False
            
                # 0. Capture standalone PICTURE shapes for quiztime layouts
                if layout.name.startswith('LAYOUT_sst_quiztime_page') and shape.shape_type == 13:  # PICTURE
                    templates['picture'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                    continue

                # 1. Identify what this shape is
                if shape.shape_type == 6:  # GROUP
                    for sub in shape.shapes:
                        if getattr(sub, 'has_text_frame', False):
                            txt_low = sub.text.lower()
                            cleaned_sub_txt = clean(sub.text)
                            if 'subtopic' in cleaned_sub_txt: has_subtopic = True
                            elif 'topic' in cleaned_sub_txt: has_topic = True
                            elif 'quiz time' in txt_low: has_quiz_title = True
                            elif 'text goes here' in txt_low or 'question goes here' in txt_low: has_body_text = True
                elif getattr(shape, 'has_text_frame', False):
                    txt_low = shape.text.lower()
                    cleaned_shape_txt = clean(shape.text)
                    if 'subtopic' in cleaned_shape_txt: has_subtopic = True
                    elif 'topic' in cleaned_shape_txt: has_topic = True
                    elif 'quiz time' in txt_low: has_quiz_title = True
                    elif 'text goes here' in txt_low or 'question goes here' in txt_low or \
                         'click to edit master' in txt_low or \
                         (shape.is_placeholder and shape.placeholder_format.type == 2) or \
                         (layout.name == 'LAYOUT_sst_discussion_page' and 'question1' in txt_low):
                        has_body_text = True
            
                # 2. Capture and remove based on identification
                if has_quiz_title:
                    templates['title'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                elif has_topic:
                    templates['topic'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                elif has_subtopic:
                    templates['subtopic'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                elif has_body_text:
                    if layout.name.startswith('LAYOUT_sst_quiztime_page'):
                        txt = shape.text.lower()
                        if 'question' in txt: templates['question'] = copy.deepcopy(shape.element)
                        elif 'options' in txt: templates['options'].append(copy.deepcopy(shape.element))
                        else: 
                            if not templates['title']: templates['title'] = copy.deepcopy(shape.element)
                    elif layout.name == 'LAYOUT_sst_discussion_page':
                        templates['question1'] = copy.deepcopy(shape.element)
                    else:
                        if isinstance(templates.get('text'), list):
                            templates['text'].append(copy.deepcopy(shape.element))
                        else:
                            templates['text'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                else:
                    # 3. If not a placeholder, check if it should be a static element
                    # Skip Picture Placeholders (Type 18)
                    if getattr(shape, 'is_placeholder', False) and shape.placeholder_format.type == 18:
                        continue
                        
                    is_static_layout = layout.name in (
                        'LAYOUT_sst_discussion_page', 'LAYOUT_sst_homework_page', 
                        'LAYOUT_syr', 'LAYOUT_ask_question', 'LAYOUT_sst_activity_page_01', 
                        'LAYOUT_sst_activity_page_02', 'LAYOUT_sst_content_page_01', 
                        '1_LAYOUT_sst_content_page_01', 'LAYOUT_sst_content_page_02', 
                        'LAYOUT_sst_deafult_page', 'LAYOUT_math_default_page',
                        'LAYOUT_sst_notedown_page', 'LAYOUT_sst_activity_static_page'
                    )
                    if is_static_layout:
                        templates['static_elements'].append(copy.deepcopy(shape.element))
                        shape.element.getparent().remove(shape.element)

    def iter_block_items(parent):
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            raise ValueError("iter_block_items: parent must be Document or _Cell")

        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)


    # Parse DOCX into sections
    sections = []
    current_section = None

    w_ns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'

    for item in iter_block_items(doc):
        lines_to_process = []
        images_to_add = []
    
        if isinstance(item, Paragraph):
            # Detect indent level from numPr/ilvl
            ilvl = 0
            is_list = False
            numPr = item._element.find(f'.//{{{w_ns}}}numPr')
            if numPr is not None:
                is_list = True
                ilvl_elem = numPr.find(f'{{{w_ns}}}ilvl')
                if ilvl_elem is not None:
                    ilvl = int(ilvl_elem.get(f'{{{w_ns}}}val', '0'))
            
            # Also check the style name for list level (handles 'List Bullet 2', etc.)
            # This serves as both detection and level override when ilvl is not set by numPr
            style_name = item.style.name if item.style else ''
            if style_name and ('List' in style_name or 'Bullet' in style_name):
                is_list = True
                import re as _re
                m = _re.search(r'\d+', style_name)
                if m:
                    style_level = max(0, int(m.group(0)) - 1)
                    if ilvl == 0 and style_level > 0:
                        ilvl = style_level
        
            # Check for math equations (OMML) and normal text runs
            # We iterate through all children of the paragraph element to preserve order
            m_ns = {'m': 'http://schemas.openxmlformats.org/officeDocument/2006/math'}
            parts = []
            
            for child in item._element.iterchildren():
                if child.tag.endswith('}oMath'):
                    from lxml import etree
                    math_xml = etree.tostring(child, encoding='unicode')
                    parts.append({'type': 'math', 'value': math_xml})
                elif child.tag.endswith('}oMathPara'):
                    # oMathPara contains one or more oMath elements
                    for subchild in child.iterchildren():
                        if subchild.tag.endswith('}oMath'):
                            from lxml import etree
                            math_xml = etree.tostring(subchild, encoding='unicode')
                            parts.append({'type': 'math', 'value': math_xml})
                elif child.tag.endswith('}r'):
                    t_elem = child.find(f'.//{{{w_ns}}}t')
                    if t_elem is not None and t_elem.text:
                        # Check formatting properties (bold)
                        rPr = child.find(f'.//{{{w_ns}}}rPr')
                        is_bold = False
                        is_italic = False
                        if rPr is not None:
                            b_elem = rPr.find(f'.//{{{w_ns}}}b')
                            i_elem = rPr.find(f'.//{{{w_ns}}}i')
                            if b_elem is not None:
                                val = b_elem.get(f'{{{w_ns}}}val', '1')
                                is_bold = val not in ('0', 'false', 'False')
                            if i_elem is not None:
                                val = i_elem.get(f'{{{w_ns}}}val', '1')
                                is_italic = val not in ('0', 'false', 'False')
                                
                        parts.append({'type': 'text', 'value': t_elem.text, 'bold': is_bold, 'italic': is_italic})
        
            # If parts is empty, fall back to plain text
            if not parts:
                lines_to_process.append((item.text.strip(), ilvl, is_list))
            else:
                lines_to_process.append((parts, ilvl, is_list))
            
            # Extract images from paragraph
            for run in item.runs:
                for drawing in run._element.findall('.//w:drawing', namespaces=run._element.nsmap):
                    for blip in drawing.findall('.//a:blip', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                        embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if embed:
                            blob = doc.part.related_parts[embed].blob
                            images_to_add.append(blob)
        elif isinstance(item, Table):
            for row in item.rows:
                # Simple conversion: join first two cells if they exist with ":"
                cells = [c.text.strip() for c in row.cells]
                if len(cells) >= 2:
                    key = cells[0].rstrip(':').strip()
                    val = cells[1].lstrip(':').strip()
                    lines_to_process.append((f"{key}: {val}", 0))
                elif len(cells) == 1:
                    lines_to_process.append((cells[0], 0))

        for entry in lines_to_process:
            content_data = entry[0] if isinstance(entry, tuple) else entry
            ilvl = entry[1] if isinstance(entry, tuple) else 0
        
            if not content_data and not images_to_add:
                continue
            
            # Get plain text for regex matching
            if isinstance(content_data, list):
                search_text = "".join(p['value'] if p['type'] == 'text' else '' for p in content_data)
            else:
                search_text = str(content_data)

            # Extract the first bracketed text as the layout name (e.g., [LAYOUT_sst_content_page_01])
            match = re.search(r'^\[\s*([^\]]+?)\s*\]', search_text)
            if match:
                # Check if there's a SYR marker on the SAME line as the layout marker
                lower_text = search_text.lower()
                has_syr = '[add_syr]' in lower_text or '[syr]' in lower_text
            
                # Check for ask_question marker on the SAME line
                aq_match = re.search(r'\[(?:add_question|ask_question)\s*\((.*?)\)\]', search_text, flags=re.IGNORECASE)
                ask_question_text = aq_match.group(1).strip() if aq_match else None
            
                name = match.group(1).strip()
            
                current_section = {
                    'name': name,
                    'content': [],
                    'images': [],
                    'has_syr': has_syr,
                    'ask_question_text': ask_question_text
                }
                sections.append(current_section)
            elif current_section is not None:
                if content_data:
                    # Use search_text for marker detection/removal
                    lower_text = search_text.lower()
                    if '[add_syr]' in lower_text or '[syr]' in lower_text:
                        current_section['has_syr'] = True
                        # If it's a string, remove it. If it's a list, it's complex, 
                        # but usually markers are at the start/end of a paragraph.
                        if isinstance(content_data, str):
                            content_data = re.sub(r'\[add_syr\]|\[syr\]', '', content_data, flags=re.IGNORECASE).strip()
                
                    # Check for ask_question marker
                    aq_match = re.search(r'\[(?:add_question|ask_question)\s*\((.*?)\)\]', search_text, flags=re.IGNORECASE)
                    if aq_match:
                        current_section['ask_question_text'] = aq_match.group(1).strip()
                        if isinstance(content_data, str):
                            content_data = re.sub(r'\[(?:add_question|ask_question)\s*\(.*?\)\]', '', content_data, flags=re.IGNORECASE).strip()

                    if content_data:
                        current_section['content'].append((content_data, ilvl))
                if images_to_add:
                    current_section['images'].extend(images_to_add)
                    images_to_add = []

    report_progress(10, "Document parsed. Generating slides...")

    def replace_text_preserve_format(shape, new_text, center=False, font_color=None, layout_name=None, is_body_text=False):
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        if not tf.paragraphs or not tf.paragraphs[0].runs:
            if isinstance(new_text, list):
                shape.text = "\n".join([str(t[0] if isinstance(t, tuple) else t) for t in new_text])
            else:
                shape.text = str(new_text)
            if center:
                for p in tf.paragraphs: p.alignment = PP_ALIGN.CENTER
            return
        
        p0 = tf.paragraphs[0]
        p0_xml = copy.deepcopy(p0._p)
    
        font = p0.runs[0].font
        name = font.name
        size = font.size
        bold = font.bold
        italic = font.italic
        underline = font.underline
    
        try:
            color_type = getattr(font.color, 'type', None)
            color_rgb = getattr(font.color, 'rgb', None) if color_type == 1 else None
            color_theme = getattr(font.color, 'theme_color', None) if color_type == 2 else None
        except:
            color_type, color_rgb, color_theme = None, None, None
    
        tf.clear()
    
        texts = new_text if isinstance(new_text, list) else [new_text]
    
        for i, para_data in enumerate(texts):
            if i == 0:
                p = tf.paragraphs[0]
                new_p_xml = copy.deepcopy(p0_xml)
                p._p.getparent().replace(p._p, new_p_xml)
                p = tf.paragraphs[0]
            else:
                new_p_xml = copy.deepcopy(p0_xml)
                tf._txBody.append(new_p_xml)
                p = tf.paragraphs[i]
                p.space_before = Pt(12)
            
            # para_data can be a string, a tuple (content, level), or a list of parts
            level = 0
            is_list = False
            if isinstance(para_data, tuple):
                level = para_data[1]
                if len(para_data) > 2:
                    is_list = para_data[2]
                para_data = para_data[0]

            # Set the paragraph indent level
            p.level = level

            # Bullet handling:
            a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            pPr = p._p.find(f'{{{a_ns}}}pPr')
            if pPr is None:
                from lxml.etree import SubElement
                pPr = SubElement(p._p, f'{{{a_ns}}}pPr')
                p._p.insert(0, pPr)

            is_sst_content = is_body_text and layout_name in ('LAYOUT_sst_content_page_01', '1_LAYOUT_sst_content_page_01', 'LAYOUT_sst_content_page_02')
            
            if is_sst_content:
                # User wants any paragraph in sst_content_page body to be a bullet
                is_list = True

                if True:
                    from lxml.etree import fromstring as parse_xml
                    # Clean existing bullet definitions
                    for child in list(pPr):
                        tag = child.tag
                        if any(tag.endswith(s) for s in ('}buNone', '}buAutoNum', '}buChar', '}buFont', '}buClr', '}buSzPct', '}buSzPts', '}buClrTx', '}buFontTx')):
                            pPr.remove(child)

                    # Determine bullet character based on level
                    if level == 0:
                        char = "❖"
                        typeface = "Noto Sans Symbols"
                    elif level == 1:
                        char = "-"
                        typeface = "Arial"
                    elif level == 2:
                        char = "▪"
                        typeface = "Arial"
                    else:
                        char = "•"
                        typeface = "Arial"

                    pPr.insert(0, parse_xml(f'<a:buChar char="{char}" xmlns:a="{a_ns}"/>'))
                    pPr.insert(0, parse_xml(f'<a:buFont typeface="{typeface}" xmlns:a="{a_ns}"/>'))
                    
                    # Set proper hanging indent using margins
                    if level == 0:
                        indent_val = 731520  # 0.8 inches padding for level 0
                        marL_val = 731520
                    else:
                        indent_val = 457200  # 0.5 inches padding for sub-bullets
                        marL_val = 731520 + (level * 457200)

                    pPr.set('marL', str(marL_val))
                    pPr.set('indent', str(-indent_val))


            # Aggressively clear existing runs and fields in the paragraph XML
            for r_elem in p._p.findall('.//a:r', namespaces=p._p.nsmap):
                p._p.remove(r_elem)
            for fld_elem in p._p.findall('.//a:fld', namespaces=p._p.nsmap):
                p._p.remove(fld_elem)
            
            if center:
                p.alignment = PP_ALIGN.CENTER
            
            # parts can be a list of math/text chunks OR a single string
            parts = para_data if isinstance(para_data, list) else [{'type': 'text', 'value': str(para_data)}]
            # Layouts that should have bold text colored yellow (#FFC000)
            yellow_bold_layouts = [
                'LAYOUT_sst_lo_page', 'LAYOUT_sst_summary_page', 'LAYOUT_sst_previous_page',
                'LAYOUT_math_lo_page', 'LAYOUT_math_summary_page'
            ]
            
            # Layouts that should have bold text colored cyan (#00FFFF)
            cyan_bold_layouts = [
                'LAYOUT_sst_content_page_01', '1_LAYOUT_sst_content_page_01', 'LAYOUT_sst_content_page_02'
            ]

            for part in parts:
                if part['type'] == 'text':
                    new_run = p.add_run()
                    new_run.text = part['value']
                
                    if name is not None: new_run.font.name = name
                    if size is not None: new_run.font.size = size
                    
                    # Apply bold/italic if explicitly set in the part, otherwise fallback to template font style
                    part_bold = part.get('bold', False)
                    if part_bold or bold is not None: 
                        new_run.font.bold = part_bold or bold
                    
                    part_italic = part.get('italic', False)
                    if part_italic or italic is not None: 
                        new_run.font.italic = part_italic or italic
                        
                    if underline is not None: new_run.font.underline = underline
                
                    # Color logic: overrides layout specific bold color, then parameter font_color, then template
                    if part_bold and layout_name in yellow_bold_layouts:
                        new_run.font.color.rgb = RGBColor(255, 192, 0) # #FFC000
                    elif part_bold and layout_name in cyan_bold_layouts:
                        new_run.font.color.rgb = RGBColor(0, 255, 255) # #00FFFF
                    elif font_color is not None:
                        new_run.font.color.rgb = font_color
                    else:
                        if color_rgb is not None:
                            new_run.font.color.rgb = color_rgb
                        elif color_theme is not None:
                            new_run.font.color.theme_color = color_theme
                elif part['type'] == 'math':
                    # Inject OMML XML
                    try:
                        from lxml.etree import fromstring as parse_xml
                        math_elem = parse_xml(part['value'])
                        p._p.append(math_elem)
                    except Exception as e:
                        print(f"Error injecting math XML: {e}")

    def apply_metadata_to_slide(slide, slide_data):
        """Recursively process all shapes in a slide and update metadata placeholders."""
        # Merge local slide_data with global_metadata, local taking precedence
        merged_data = {**global_metadata, **slide_data}
        
        def process_shape_list(shapes, parent_group=None):
            for shape in shapes:
                # print(f"DEBUG: Processing shape {shape.name if hasattr(shape, 'name') else shape} type={shape.shape_type}", flush=True)
                if shape.shape_type == 6:  # msoGroup
                    process_shape_list(shape.shapes, parent_group=shape)
                elif getattr(shape, 'has_text_frame', False):
                    cleaned_txt = clean(shape.text)
                    if not cleaned_txt:
                        continue
                        
                    # Priority: Subtopic (must check before topic because 'topic' is in 'subtopic')
                    if 'subtopic' in cleaned_txt:
                        txt_low = shape.text.lower()
                        if 'topic' in txt_low.replace('subtopic', ''):
                            # Combined placeholder: Topic Name - Subtopic Name
                            topic_val = merged_data.get('topic', '')
                            subtopic_val = merged_data.get('subtopic', '')
                            if topic_val and subtopic_val:
                                val = f"{topic_val} - {subtopic_val}"
                            else:
                                val = topic_val or subtopic_val
                        else:
                            val = merged_data.get('subtopic')
                            
                        if val:
                            replace_text_preserve_format(shape, val, center=True, layout_name=slide.slide_layout.name)
                            tf = shape.text_frame
                            tf.word_wrap = False
                            # Estimate required width (250k EMUs per char)
                            # Add just enough padding for the diagonal corners
                            text_width = len(str(val).strip()) * 250000
                            padding = 800000  # Generous padding for shape geometry
                            required_width = text_width + padding
                            
                            # Ensure we don't shrink the shape below its original template size
                            if shape.width > 0 and required_width < shape.width:
                                required_width = shape.width
                                
                            max_w = int(prs.slide_width * 0.95)
                            if required_width > max_w:
                                required_width = max_w
                                tf.word_wrap = True  # Allow wrap if it's truly massive
                            
                            # Set text frame margins to give text space within the shape
                            tf.margin_left = int(padding / 2.5)
                            tf.margin_right = int(padding / 2.5)
                            tf.margin_top = Pt(4)
                            tf.margin_bottom = Pt(4)
                            
                            if shape.width > 0:
                                scale = required_width / shape.width
                                if parent_group:
                                    # Proportional scaling of all group elements relative to group start
                                    g_left = parent_group.left
                                    for child in parent_group.shapes:
                                        child.width = int(child.width * scale)
                                        rel_left = child.left - g_left
                                        child.left = g_left + int(rel_left * scale)
                                    parent_group.width = int(parent_group.width * scale)
                                else:
                                    shape.width = required_width
                            else:
                                shape.width = required_width

                    elif 'topic' in cleaned_txt and merged_data.get('topic'):
                        replace_text_preserve_format(shape, merged_data['topic'], center=True, layout_name=slide.slide_layout.name)
                            
                    elif cleaned_txt == 'class' and merged_data.get('class'):
                        replace_text_preserve_format(shape, merged_data['class'], center=True, layout_name=slide.slide_layout.name)
                    elif cleaned_txt == 'subject' and merged_data.get('subject'):
                        replace_text_preserve_format(shape, merged_data['subject'], center=True, layout_name=slide.slide_layout.name)
                    elif cleaned_txt == 'chapternumber' and merged_data.get('chapter number'):
                        replace_text_preserve_format(shape, merged_data['chapter number'], center=True, layout_name=slide.slide_layout.name)
                    elif cleaned_txt == 'chaptername' and merged_data.get('chapter name'):
                        replace_text_preserve_format(shape, merged_data['chapter name'], center=True, layout_name=slide.slide_layout.name)
                    elif cleaned_txt == 'lesson' and merged_data.get('lesson'):
                        replace_text_preserve_format(shape, merged_data['lesson'], center=True, layout_name=slide.slide_layout.name)

        process_shape_list(slide.shapes)

    total_sections = len(sections)
    for i, section in enumerate(sections):
        progress = 10 + int((i / total_sections) * 85)
        report_progress(progress, f"Generating slide {i+1} of {total_sections}: {section['name']}...")
        
        sname = section['name'].strip().lower()
        if sname == 'sst_content_page':
            if len(section.get('images', [])) > 1:
                layout = get_layout('sst_content_page_02')
            else:
                layout = get_layout('sst_content_page_01')
        elif sname in ('sst_quiztime_page', 'quiztime_page', 'sstquiztimepage'):
            # Parse quiz content first to determine layout
            quiz_data = {'question': '', 'options': []}
            for entry in section['content']:
                line = get_text(entry)
                line_low = line.lower()
                if 'question:' in line_low:
                    quiz_data['question'] = line.split(':', 1)[1].strip()
                elif 'options:' in line_low:
                    opts = line.split(':', 1)[1].strip()
                    quiz_data['options'] = [o.strip() for o in opts.split(',')]
                elif line.strip():
                    if not quiz_data['question']: quiz_data['question'] = line.strip()
                    else: quiz_data['options'].append(line.strip())
        
            # Pick layout based on option character length
            # If ANY option > 25 chars → _01 (long), else _02 (short)
            use_long = any(len(opt) > 25 for opt in quiz_data['options'])
            if use_long:
                layout = get_layout('sst_quiztime_page_01')
                print(f"Quiz options have long text (>25 chars) -> using LAYOUT_sst_quiztime_page_01")
            else:
                layout = get_layout('sst_quiztime_page_02')
                print(f"Quiz options are short (<=25 chars) -> using LAYOUT_sst_quiztime_page_02")
        elif sname in ('sst_activity_page', 'activity_page', 'activity'):
            if len(section.get('images', [])) > 1:
                layout = get_layout('sst_activity_page_01')
            elif len(section['content']) > 1:
                layout = get_layout('sst_activity_page_02')
            else:
                layout = get_layout('sst_activity_page_01')
        else:
            layout = get_layout(section['name'])

        if not layout:
            print(f"Skipping section [{section['name']}], layout not found.")
            continue
        
        slide = prs.slides.add_slide(layout)
        
        # Inject the captured logo elements to every slide
        # inject_logo(slide)
        if section['name'].replace("_", "").lower() in ('mathpagetitle', 'mathtitlepage', 'sstpagetitle', 'ssttitlepage'):
            data = {}
            for entry in section['content']:
                line = get_text(entry)
                if ":" in line:
                    parts = line.split(":", 1)
                    key = clean(parts[0])
                    val = parts[1].strip()
                    data[key.upper()] = val
                    # Store in global metadata if key matches
                    for g_key in global_metadata:
                        if clean(g_key) == key:
                            global_metadata[g_key] = val
                            break
        
            idx_mapping = {}
            for shape in layout.shapes:
                if shape.is_placeholder and shape.has_text_frame:
                    idx_mapping[clean(shape.text)] = shape.placeholder_format.idx
                
            mapping = {
                "CLASS": "class",
                "SUBJECT": "subject",
                "CHAPTER_NUMBER": "chapter number",
                "CHAPTER_NAME": "chapter name",
                "LESSON": "lesson",
                "TOPIC": "topic"
            }

            for key, template_word in mapping.items():
                cleaned_word = clean(template_word)
                val = data.get(key) or global_metadata.get(template_word)
                if val:
                    if cleaned_word in idx_mapping:
                        idx = idx_mapping[cleaned_word]
                        found_on_slide = False
                        for shape in slide.shapes:
                            if shape.is_placeholder and shape.placeholder_format.idx == idx:
                                shape.text = val
                                print(f"Updated PageTitle: {template_word} -> {val}")
                                found_on_slide = True
                                break
                        if not found_on_slide:
                            print(f"[{template_word}] idx {idx} NOT FOUND on slide shapes!")
                    else:
                        print(f"[{template_word}] NOT FOUND in idx_mapping! keys={idx_mapping.keys()}")
                else:
                    print(f"[{key}] NOT FOUND in docx data or global metadata!")
            
            # Catch-all for any other metadata placeholders on title slide
            apply_metadata_to_slide(slide, data)
        elif layout.name in lo_group_xmls:
            # Inject static title if available
            if layout.name in lo_title_xmls and lo_title_xmls[layout.name] is not None:
                slide.shapes._spTree.append(copy.deepcopy(lo_title_xmls[layout.name]))

            group_xml = lo_group_xmls.get(layout.name)
            
            if group_xml is not None:
                # Group content: merge sub-bullets (ilvl > 0) with their parent line
                # Each item: (main_text, [(sub_text, ilvl), ...])
                grouped_items = []
                for entry in section['content']:
                    raw_text = entry[0] if isinstance(entry, tuple) else entry
                    ilvl = entry[1] if isinstance(entry, tuple) else 0
                    
                    # Convert raw_text to string for checking
                    if isinstance(raw_text, list):
                        plain_line = "".join([p['value'] for p in raw_text if p['type'] == 'text']).strip()
                    else:
                        plain_line = str(raw_text).strip()
                    
                    if plain_line.startswith(('•', '-', '*')):
                        plain_line = plain_line.lstrip('•-*').strip()
                        # If raw_text is a list, we should technically trim the first text part. 
                        # For simplicity, we just keep raw_text as is but use plain_line for grouping logic.
                        # Wait, we need to strip the bullet from the actual printed text too.
                        if isinstance(raw_text, list):
                            if raw_text and raw_text[0]['type'] == 'text':
                                raw_text[0]['value'] = raw_text[0]['value'].lstrip('•-*').strip()
                        else:
                            raw_text = str(raw_text).lstrip('•-*').strip()
                
                    if ilvl == 0:
                        grouped_items.append((raw_text, []))
                    elif ilvl > 0 and grouped_items:
                        grouped_items[-1][1].append((raw_text, ilvl))
                    else:
                        grouped_items.append((raw_text, []))
            
                current_top_offset = 0
                for main_text, sub_entries in grouped_items:
                    new_element = copy.deepcopy(group_xml)
                    slide.shapes._spTree.append(new_element)
                    new_shape = slide.shapes[-1]
                    new_shape.top = new_shape.top + current_top_offset
                
                    for subshape in new_shape.shapes:
                        if getattr(subshape, 'has_text_frame', False) and 'Text goes here' in subshape.text:
                            tf = subshape.text_frame
                            a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                            
                            # Remove all paragraphs except the first (template subtitle paras)
                            txBody = tf._txBody
                            all_paras = txBody.findall(f'{{{a_ns}}}p')
                            for extra_p in all_paras[1:]:
                                txBody.remove(extra_p)
                            
                            # Vertically align text at top to match decorative shape position
                            bodyPr = txBody.find(f'{{{a_ns}}}bodyPr')
                            if bodyPr is not None:
                                bodyPr.set('anchor', 't')
                            
                            # Set main text in first paragraph
                            p0 = tf.paragraphs[0]
                            # Clear existing runs
                            for r_len in range(len(p0.runs)-1, -1, -1):
                                r_elem = p0.runs[r_len]._r
                                r_elem.getparent().remove(r_elem)
                                
                            parts = main_text if isinstance(main_text, list) else [{'type': 'text', 'value': str(main_text)}]
                            for part in parts:
                                if part['type'] == 'text':
                                    run = p0.add_run()
                                    run.text = part['value']
                                    run.font.size = Pt(36)
                                    part_bold = part.get('bold', False)
                                    if part_bold and layout.name in ['LAYOUT_sst_lo_page', 'LAYOUT_sst_summary_page', 'LAYOUT_sst_previous_page', 'LAYOUT_math_lo_page', 'LAYOUT_math_summary_page']:
                                        run.font.bold = part_bold
                                        run.font.color.rgb = RGBColor(255, 192, 0)
                                    else:
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                            
                            # Set spacing on main paragraph (0 before to align with shape, 600 after for gap)
                            p0_pPr = tf.paragraphs[0]._p.find(f'{{{a_ns}}}pPr')
                            if p0_pPr is not None:
                                spcBef = p0_pPr.find(f'{{{a_ns}}}spcBef')
                                if spcBef is not None:
                                    spcPts = spcBef.find(f'{{{a_ns}}}spcPts')
                                    if spcPts is not None:
                                        spcPts.set('val', '0')  # No space before main text
                                spcAft = p0_pPr.find(f'{{{a_ns}}}spcAft')
                                if spcAft is not None:
                                    spcPts = spcAft.find(f'{{{a_ns}}}spcPts')
                                    if spcPts is not None:
                                        spcPts.set('val', '600')
                        
                            # Add sub-bullet paragraphs using cloned template XML
                            subtitle_templates = lo_subtitle_para_xmls.get(layout.name, [])
                            for bullet_text, bullet_ilvl in sub_entries:
                                # Pick the right template: ilvl 1 -> index 0, ilvl 2 -> index 1, etc.
                                tmpl_idx = min(bullet_ilvl - 1, len(subtitle_templates) - 1) if subtitle_templates else -1
                                if tmpl_idx >= 0:
                                    # Clone the subtitle paragraph template
                                    new_p = copy.deepcopy(subtitle_templates[tmpl_idx])
                                    # Clear existing text runs
                                    for r_elem in new_p.findall(f'{{{a_ns}}}r'):
                                        new_p.remove(r_elem)
                                    for fld_elem in new_p.findall(f'{{{a_ns}}}fld'):
                                        new_p.remove(fld_elem)
                                    # Remove endParaRPr if exists, we'll add fresh run
                                    for endPr in new_p.findall(f'{{{a_ns}}}endParaRPr'):
                                        new_p.remove(endPr)
                                    
                                    # Set consistent spacing
                                    pPr = new_p.find(f'{{{a_ns}}}pPr')
                                    if pPr is not None:
                                        spcBef = pPr.find(f'{{{a_ns}}}spcBef')
                                        if spcBef is not None:
                                            spcPts = spcBef.find(f'{{{a_ns}}}spcPts')
                                            if spcPts is not None:
                                                spcPts.set('val', '600')  # 6pt consistent spacing
                                        spcAft = pPr.find(f'{{{a_ns}}}spcAft')
                                        if spcAft is not None:
                                            spcPts = spcAft.find(f'{{{a_ns}}}spcPts')
                                            if spcPts is not None:
                                                spcPts.set('val', '600')
                                    
                                    parts = bullet_text if isinstance(bullet_text, list) else [{'type': 'text', 'value': str(bullet_text)}]
                                    for part in parts:
                                        if part['type'] == 'text':
                                            color_val = "FFFFFF"  # White default
                                            part_bold = part.get('bold', False)
                                            # If part is bold and on a target layout, make it yellow (#FFC000)
                                            if part_bold and layout.name in ['LAYOUT_sst_lo_page', 'LAYOUT_sst_summary_page', 'LAYOUT_sst_previous_page', 'LAYOUT_math_lo_page', 'LAYOUT_math_summary_page']:
                                                color_val = "FFC000"
                                            
                                            bold_attr = ' b="1"' if part_bold else ''
                                            
                                            r_elem = parse_xml(
                                                f'<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                                                f'<a:rPr lang="en-US" sz="{max(2000, 3200 - (bullet_ilvl * 400))}" dirty="0"{bold_attr}>'
                                                f'<a:solidFill><a:srgbClr val="{color_val}"/></a:solidFill>'
                                                f'<a:latin typeface="Calibri"/>'
                                                f'<a:ea typeface="Calibri"/>'
                                                f'<a:cs typeface="Calibri"/>'
                                                f'</a:rPr>'
                                                f'<a:t>{part["value"]}</a:t>'
                                                f'</a:r>'
                                            )
                                            new_p.append(r_elem)
                                    txBody.append(new_p)
                                else:
                                    # Fallback if no subtitle templates are available
                                    p = tf.add_paragraph()
                                    p.text = bullet_text
                                    p.space_before = Pt(6)
                                    p.space_after = Pt(6)
                                    p.level = bullet_ilvl
                                    for run in p.runs:
                                        run.font.size = Pt(max(20, 36 - (bullet_ilvl * 4)))
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                
                    # Calculate height based on content
                    chars_per_line = 65
                    line_count = max(1, len(main_text) // chars_per_line + (1 if len(main_text) % chars_per_line > 0 else 0))
                    sub_height = len(sub_entries) * 0.45
                    current_top_offset += int(Inches(0.9 + 0.55 * (line_count - 1) + sub_height))
                
                total_subs = sum(len(s) for _, s in grouped_items)
                print(f"Updated {layout.name} with {len(grouped_items)} items ({total_subs} sub-bullets)")
                
                # Use centralized metadata injection for topic/subtopic
                apply_metadata_to_slide(slide, {})
        elif layout.name in ('LAYOUT_sst_quiztime_page_01', 'LAYOUT_sst_quiztime_page_02'):
            templates = sst_content_templates.get(layout.name, {})
        
            # Inject static picture
            if templates.get('picture') is not None:
                pic_elem = copy.deepcopy(templates['picture'])
                copy_image_rels(pic_elem, layout.part, slide.part)
                slide.shapes._spTree.append(pic_elem)

            # Inject static title
            if templates.get('title') is not None:
                title_elem = copy.deepcopy(templates['title'])
                copy_image_rels(title_elem, layout.part, slide.part)
                slide.shapes._spTree.append(title_elem)

            # Reuse quiz_data if already parsed during layout selection, otherwise parse now
            if 'quiz_data' not in dir():
                quiz_data = {'question': '', 'options': []}
                for entry in section['content']:
                    content_obj = entry[0] if isinstance(entry, tuple) else entry
                    if isinstance(content_obj, list):
                        # For quiz, we usually want text but can support math in question/options
                        # For now, append the rich object to whichever field is current
                        if not quiz_data['question']: quiz_data['question'] = content_obj
                        else: quiz_data['options'].append(content_obj)
                        continue

                    line = str(content_obj)
                    line_low = line.lower()
                    if 'question:' in line_low:
                        quiz_data['question'] = line.split(':', 1)[1].strip()
                    elif 'options:' in line_low:
                        opts = line.split(':', 1)[1].strip()
                        quiz_data['options'] = [o.strip() for o in opts.split(',')]
                    elif line.strip():
                        if not quiz_data['question']: quiz_data['question'] = line.strip()
                        else: quiz_data['options'].append(line.strip())

            if templates.get('question') is not None:
                question_elem = copy.deepcopy(templates['question'])
            
                # If there are no options, vertically and horizontally center the question text
                if not quiz_data['options']:
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    body = question_elem.find('.//a:bodyPr', namespaces=ns)
                    if body is not None:
                        # Set vertical alignment to center
                        body.set('anchor', 'ctr')

                    off = question_elem.find('.//a:xfrm/a:off', namespaces=ns)
                    ext = question_elem.find('.//a:xfrm/a:ext', namespaces=ns)
                    if off is not None and ext is not None:
                        slide_height = prs.slide_height
                        slide_width = prs.slide_width
                        shape_height = int(ext.get('cy', '0'))
                        shape_width = int(ext.get('cx', '0'))
                    
                        # Center shape vertically (assuming we want to center the whole box and its text)
                        # Center Y: taking into account the space available below the title
                        # But simpler is just center vertically on slide
                        new_y = int((slide_height - shape_height) / 2)
                        
                        # Optionally center horizontally (already done by the previous code block)
                        new_x = int((slide_width - shape_width) / 2)
                        off.set('y', str(new_y))
                        off.set('x', str(new_x))
            
                slide.shapes._spTree.append(question_elem)
                # Pass center=True to replace_text_preserve_format if there are no options to center the paragraph text
                replace_text_preserve_format(slide.shapes[-1], quiz_data['question'], center=not bool(quiz_data['options']))
        
            # Inject options only if they exist
            if quiz_data['options']:
                option_templates = templates.get('options', [])
                if len(option_templates) == 1:
                    # _01 style: single placeholder, all options as multi-line text
                    slide.shapes._spTree.append(copy.deepcopy(option_templates[0]))
                    replace_text_preserve_format(slide.shapes[-1], quiz_data['options'])
                elif len(option_templates) > 1:
                    # _02 style: separate placeholder per option
                    for idx, opt_template in enumerate(option_templates):
                        slide.shapes._spTree.append(copy.deepcopy(opt_template))
                        if idx < len(quiz_data['options']):
                            replace_text_preserve_format(slide.shapes[-1], quiz_data['options'][idx])
                        else:
                            replace_text_preserve_format(slide.shapes[-1], '')
        
            print(f"Updated {layout.name} with quiz question and {len(quiz_data['options'])} options.")
            
            # Use centralized metadata injection for topic/subtopic
            apply_metadata_to_slide(slide, {})
            quiz_data = None  # Reset for next quiz section
        elif layout.name in ('LAYOUT_sst_content_page_01', 'LAYOUT_sst_content_page_02', 'LAYOUT_sst_deafult_page', 'LAYOUT_math_default_page'):
            data = {}
            data_text_list = []
            has_local_topic = False
            for entry in section['content']:
                content_obj = entry[0] if isinstance(entry, tuple) else entry
                ilvl = entry[1] if isinstance(entry, tuple) else 0
                is_list = entry[2] if (isinstance(entry, tuple) and len(entry) > 2) else False
            
                if isinstance(content_obj, list):
                    # For parsing Topic/Subtopic, we need the plain text
                    line = "".join([p['value'] for p in content_obj if p['type'] == 'text'])
                else:
                    line = str(content_obj)
                    
                line_low = line.lower()
            
                if line_low.startswith('topic:'):
                    data['topic'] = line.split(':', 1)[1].strip()
                    if data['topic']: has_local_topic = True
                elif line_low.startswith('subtopic:'):
                    data['subtopic'] = line.split(':', 1)[1].strip()
                elif ":" in line:
                    parts = line.split(":", 1)
                    key = clean(parts[0])
                    val = parts[1].strip()
                    if 'subtopic' in key or key.startswith('subtopic') or key.startswith('subtop'):
                        data['subtopic'] = val
                    elif 'topic' in key or key.startswith('topic'):
                        data['topic'] = val
                    elif key == 'text':
                        data_text_list.append((content_obj, ilvl, is_list))
                    else:
                        data_text_list.append((content_obj, ilvl, is_list))
                elif line.strip():
                    data_text_list.append((content_obj, ilvl, is_list))
                
            # Fallback to global metadata
            if 'topic' not in data and global_metadata.get('topic'):
                data['topic'] = global_metadata['topic']
            if 'subtopic' not in data and global_metadata.get('subtopic'):
                data['subtopic'] = global_metadata['subtopic']

            if data_text_list:
                data['text'] = data_text_list
        
            # Inject extracted XML templates onto this slide from the correct layout version
            shape_elements = [] # All elements injected
            text_shapes = []    # Only elements that should hold body text
            templates = sst_content_templates.get(layout.name, {})
        
            # Inject topic shape if available (local or global)
            if templates.get('topic') is not None and 'topic' in data:
                topic_elem = copy.deepcopy(templates['topic'])
                slide.shapes._spTree.append(topic_elem)
                shape_elements.append(slide.shapes[-1])
            
            if templates.get('subtopic') is not None and data.get('subtopic'):
                subtopic_elem = copy.deepcopy(templates['subtopic'])
                slide.shapes._spTree.append(subtopic_elem)
                shape_elements.append(slide.shapes[-1])
            
            if templates.get('static_elements'):
                for static_elem in templates['static_elements']:
                    slide.shapes._spTree.append(copy.deepcopy(static_elem))

            if templates.get('text') is not None and len(templates['text']) > 0:
                texts = templates['text'] if isinstance(templates['text'], list) else [templates['text']]
                for text_elem_xml in texts:
                    text_elem = copy.deepcopy(text_elem_xml)
                    
                    # If no subtopic in data, shift the text box up to the subtopic's position
                    if not data.get('subtopic') and templates.get('subtopic') is not None:
                        sub_xml = templates['subtopic']
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        sub_off = sub_xml.find('.//a:off', namespaces=ns)
                        text_off = text_elem.find('.//a:off', namespaces=ns)
                        if sub_off is not None and text_off is not None:
                            text_off.set('y', sub_off.get('y'))
                
                    slide.shapes._spTree.append(text_elem)
                    new_shape = slide.shapes[-1]
                    shape_elements.append(new_shape)
                    
                    # Identify if this injected shape is a text shape
                    if getattr(new_shape, 'shape_type', None) == 6: # Group
                        for sub in new_shape.shapes:
                            if getattr(sub, 'has_text_frame', False) and \
                               ('text goes here' in sub.text.lower() or 'click to edit' in sub.text.lower()):
                                text_shapes.append(sub)
                    elif getattr(new_shape, 'has_text_frame', False):
                        # Even if it contains other master text, if it came from the 'text' template, it's our target
                        text_shapes.append(new_shape)
            
            # Insert text and subtexts
            if text_shapes and 'text' in data:
                paragraphs = data['text']
                if len(text_shapes) > 1 and len(paragraphs) > 1:
                    # Distribute paragraphs across text shapes evenly
                    from math import ceil
                    chunk_size = ceil(len(paragraphs) / len(text_shapes))
                    for idx, shape in enumerate(text_shapes):
                        chunk = paragraphs[idx * chunk_size : (idx + 1) * chunk_size]
                        if chunk:
                            replace_text_preserve_format(shape, chunk, layout_name=layout.name, is_body_text=True)
                        else:
                            replace_text_preserve_format(shape, "", layout_name=layout.name, is_body_text=True)
                else:
                    # Only one text shape or one paragraph, put everything in the first shape and clear others
                    replace_text_preserve_format(text_shapes[0], paragraphs, layout_name=layout.name, is_body_text=True)
                    for shape in text_shapes[1:]:
                        replace_text_preserve_format(shape, "", layout_name=layout.name, is_body_text=True)
            elif 'text' in data:
                pass # No matching placeholders found for the provided text

            # Use centralized metadata injection for topic/subtopic
            apply_metadata_to_slide(slide, data)
                            
            # Handle picture placeholder(s)
            images = section.get('images', [])
            pic_ph = None
            for shape in slide.shapes:
                if getattr(shape, 'is_placeholder', False) and shape.placeholder_format.type == 18:
                    pic_ph = shape
                    break
        
            if pic_ph is not None and images:
                # Capture properties of the placeholder before it's potentially replaced/removed
                left, top, width, height = pic_ph.left, pic_ph.top, pic_ph.width, pic_ph.height
                num_images = len(images)
            
                if num_images > 1:
                    # Remove the placeholder to avoid overlap with newly added pictures
                    pic_ph.element.getparent().remove(pic_ph.element)
                
                    gap = Inches(0.1)
                    total_gap = gap * (num_images - 1)
                    img_width = (width - total_gap) / num_images
                
                    for i, img_blob in enumerate(images):
                        # Load image to get original dimensions for aspect ratio
                        with Image.open(io.BytesIO(img_blob)) as pil_img:
                            orig_w, orig_h = pil_img.size
                    
                        slot_w = img_width
                        slot_h = height
                    
                        # Calculate fit dimensions
                        slot_aspect = slot_w / slot_h
                        img_aspect = orig_w / orig_h
                    
                        if img_aspect > slot_aspect:
                            # Image is wider than slot relative to height
                            draw_w = slot_w
                            draw_h = slot_w / img_aspect
                        else:
                            # Image is taller than slot relative to width
                            draw_h = slot_h
                            draw_w = slot_h * img_aspect
                    
                        # Center within slot
                        final_left = int(left + (i * (slot_w + gap)) + (slot_w - draw_w) / 2)
                        final_top = int(top + (slot_h - draw_h) / 2)
                    
                        slide.shapes.add_picture(
                            io.BytesIO(img_blob), 
                            final_left, 
                            final_top, 
                            width=int(draw_w),
                            height=int(draw_h)
                        )
                else:
                    # Single image: use standard placeholder insertion (handles aspect ratio better)
                    pic_ph.insert_picture(io.BytesIO(images[0]))
            elif pic_ph is not None:
                # Hide/Remove placeholder if no images are present
                pic_ph.element.getparent().remove(pic_ph.element)
                    
            print(f"Updated {layout.name} with text and {len(images)} images.")

        elif layout.name in ('LAYOUT_sst_activity_page_01', 'LAYOUT_sst_activity_page_02'):
            templates = sst_content_templates.get(layout.name, {})
        
            # Inject static elements
            for static_elem in templates.get('static_elements', []):
                elem_copy = remove_locks(copy.deepcopy(static_elem))
                copy_image_rels(elem_copy, layout.part, slide.part)
                slide.shapes._spTree.append(elem_copy)

            # Handle text injection and dynamic resizing
            text_content = []
            for entry in section['content']:
                content_obj = entry[0] if isinstance(entry, tuple) else entry
                text_content.append(content_obj)

            if templates.get('text') is not None:
                text_xmls = templates['text'] if isinstance(templates['text'], list) else [templates['text']]
                if text_xmls:
                    text_elem = copy.deepcopy(text_xmls[0])
                    slide.shapes._spTree.append(text_elem)
                    shape = slide.shapes[-1]
                # Force black text for activity boxes if they are yellow
                replace_text_preserve_format(shape, text_content, font_color=RGBColor(0, 0, 0), layout_name=layout.name)
            
                # Estimate height needed for 36pt font
                chars_per_line = max(20, int((shape.width / 914400) * 3.5))
            
                line_count = 0
                for line_item in text_content:
                    # Get plain text length for estimation
                    if isinstance(line_item, list):
                        line_str = "".join(p['value'] if p['type'] == 'text' else '    ' for p in line_item)
                    else:
                        line_str = str(line_item)
                    line_count += (len(line_str) // chars_per_line) + 1
                
                if line_count > 1:
                    # 36pt font is ~0.5 inch high. With spacing, we need ~0.7 inch per line.
                    # EMU: 1 inch = 914400. 0.7 inch = ~640,000 EMUs.
                    extra_h = (line_count - 1) * 650000 
                    old_h = shape.height
                    shape.height += extra_h
                    new_h = shape.height
                
                    # Maintain corner roundness
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    gd = shape.element.find('.//a:gd[@name="adj"]', namespaces=ns)
                    if gd is not None:
                        try:
                            fmla = gd.get('fmla')
                            if 'val' in fmla:
                                old_adj = int(fmla.split()[-1])
                                new_adj = int(old_adj * (old_h / new_h))
                                gd.set('fmla', f'val {new_adj}')
                        except:
                            pass
                
                    # Move everything below it down (e.g. discussion group)
                    for other_shape in slide.shapes:
                        # Skip if it is the picture placeholder
                        if getattr(other_shape, 'is_placeholder', False) and other_shape.placeholder_format.type == 18:
                            continue
                        # Skip if it is a standard picture (the default image)
                        if other_shape.shape_type == 13: # PICTURE
                            continue
                            
                        # Use old_h to only match shapes TRULY below the text box (giving 0.1 inch buffer)
                        if other_shape != shape and other_shape.top >= (shape.top + old_h - 91440):
                             # Use a bit more buffer when moving things down
                             other_shape.top += extra_h + 100000

            # Handle picture placeholder
            images = section.get('images', [])
            pic_ph = None
            for shp in slide.shapes:
                if getattr(shp, 'is_placeholder', False) and shp.placeholder_format.type == 18:
                    pic_ph = shp
                    break
        
            if pic_ph is not None and images:
                # Remove any default images copied from the layout to avoid overlap
                for shp in list(slide.shapes):
                    if shp.shape_type == 13: # PICTURE
                        shp.element.getparent().remove(shp.element)
                        
                # Capture properties of the placeholder
                left, top, total_width, total_height = pic_ph.left, pic_ph.top, pic_ph.width, pic_ph.height
                num_images = len(images)
            
                if num_images > 1:
                    # Remove the placeholder to avoid overlap
                    pic_ph.element.getparent().remove(pic_ph.element)
                
                    # We want square slots. The width of each slot is total_width / num_images.
                    # But we should also consider the height. 
                    # Let's use the smaller of (total_width/num_images) and total_height as the square side.
                    gap = Inches(0.1)
                    slot_side = min((total_width - (gap * (num_images-1))) / num_images, total_height)
                
                    for i, img_blob in enumerate(images):
                        with Image.open(io.BytesIO(img_blob)) as pil_img:
                            orig_w, orig_h = pil_img.size
                    
                        img_aspect = orig_w / orig_h
                    
                        # Calculate dimensions to fit inside the square slot_side x slot_side
                        if img_aspect > 1: # Wide image
                            draw_w = slot_side
                            draw_h = slot_side / img_aspect
                        else: # Tall or square image
                            draw_h = slot_side
                            draw_w = slot_side * img_aspect
                        
                        # Center within the slot
                        slot_left = left + i * (slot_side + gap)
                        # For layout 01, we center vertically in the placeholder's original vertical space
                        final_left = int(slot_left + (slot_side - draw_w) / 2)
                        final_top = int(top + (total_height - draw_h) / 2)
                    
                        slide.shapes.add_picture(
                            io.BytesIO(img_blob), 
                            final_left, 
                            final_top, 
                            width=int(draw_w),
                            height=int(draw_h)
                        )
                else:
                    # Single image: use standard placeholder insertion
                    slide.shapes._spTree.append(pic_ph.element)
                    pic_ph.insert_picture(io.BytesIO(images[0]))
            elif pic_ph is not None:
                pic_ph.element.getparent().remove(pic_ph.element)

            # Use centralized metadata injection for topic/subtopic
            apply_metadata_to_slide(slide, {})
            print(f"Updated {layout.name} with activity content and {len(images)} images.")

        elif layout.name == 'LAYOUT_final_quiz_page':
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    shape.text = "Final Quiz!"
                    break
            print("Inserted LAYOUT_final_quiz_page (Welcome Slide)")

            qa_pairs = []
            for entry in section['content']:
                line = get_text(entry)
                if ":" in line:
                    parts = line.split(":", 1)
                    q_text = parts[0].strip()
                    a_text = parts[1].strip()
                    if q_text.lower() == 'question' and a_text.lower() == 'answer':
                        continue
                    qa_pairs.append((q_text, a_text))
        
            layout_q = get_layout('final_quiz_page_q')
            layout_a = get_layout('final_quiz_page_a')
        
            for q, a in qa_pairs:
                if layout_q:
                    s_q = prs.slides.add_slide(layout_q)
                    inject_logo(s_q)
                    for shape in s_q.shapes:
                        if shape.is_placeholder and shape.placeholder_format.idx == 0:
                            shape.text = q
                if layout_a:
                    s_a = prs.slides.add_slide(layout_a)
                    inject_logo(s_a)
                    for shape in s_a.shapes:
                        if shape.is_placeholder and shape.placeholder_format.idx == 0:
                            shape.text = a
            print(f"Generated {len(qa_pairs)} Question and Answer slide pairs.")
        elif layout.name == 'LAYOUT_sst_notedown_page':
            templates = sst_content_templates.get(layout.name, {})
            data = {}
            has_local_topic = False
            
            # First pass: Extract topic/subtopic and identify lines to skip
            body_entries = []
            for entry in section['content']:
                line = get_text(entry)
                line_low = line.lower()
                if line_low.startswith('topic:'):
                    data['topic'] = line.split(':', 1)[1].strip()
                    if data['topic']: has_local_topic = True
                elif line_low.startswith('subtopic:'):
                    data['subtopic'] = line.split(':', 1)[1].strip()
                else:
                    body_entries.append(entry)

            # Fallback for data only
            if 'topic' not in data and global_metadata.get('topic'):
                data['topic'] = global_metadata['topic']
            if 'subtopic' not in data and global_metadata.get('subtopic'):
                data['subtopic'] = global_metadata['subtopic']

            # 1. Inject Topic/Subtopic first so they are behind (earlier in _spTree)
            if templates.get('topic') is not None and has_local_topic:
                slide.shapes._spTree.append(copy.deepcopy(templates['topic']))
            
            if templates.get('subtopic') is not None and 'subtopic' in data:
                slide.shapes._spTree.append(copy.deepcopy(templates['subtopic']))

            # 2. Inject static elements (NoteDown header) on top of topic
            for static_elem in templates.get('static_elements', []):
                elem_copy = remove_locks(copy.deepcopy(static_elem))
                copy_image_rels(elem_copy, layout.part, slide.part)
                slide.shapes._spTree.append(elem_copy)

            # 3. Inject text box
            shape = None
            if templates.get('text') is not None:
                text_xmls = templates['text'] if isinstance(templates['text'], list) else [templates['text']]
                if text_xmls:
                    slide.shapes._spTree.append(copy.deepcopy(text_xmls[0]))
                    shape = slide.shapes[-1]
            
            if shape and body_entries:
                processed_content = []
                for entry in body_entries:
                    content_obj = entry[0] if isinstance(entry, tuple) else entry
                    ilvl = entry[1] if isinstance(entry, tuple) else 0
                    
                    # Apply diamond marker to top-level bullets
                    if ilvl == 0 and isinstance(content_obj, str):
                        clean_line = content_obj.strip()
                        if clean_line.startswith(("•", "-", "*")):
                            content_obj = "❖ " + clean_line.lstrip("•-*").strip()
                    
                    processed_content.append((content_obj, ilvl))
            
                replace_text_preserve_format(shape, processed_content, layout_name=layout.name)
                
            # Use centralized metadata injection for topic/subtopic content replacement
            apply_metadata_to_slide(slide, data)
            print(f"Updated {layout.name} with topic visibility (layered behind) and bullet support.")

        elif layout.name == 'LAYOUT_sst_discussion_page':
            templates = sst_content_templates.get(layout.name, {})
            # Inject all static elements (groups, pictures, title) from layout
            for static_elem in templates.get('static_elements', []):
                elem_copy = remove_locks(copy.deepcopy(static_elem))
                copy_image_rels(elem_copy, layout.part, slide.part)
                slide.shapes._spTree.append(elem_copy)
            if templates.get('question1') is not None:
                # Parse all content lines from the section
                q1_lines = []
                for entry in section['content']:
                    line = get_text(entry)
                    if ":" in line:
                        parts = line.split(":", 1)
                        if parts[0].strip().lower() == 'question1':
                            q1_lines.append(parts[1].strip())
                        else:
                            q1_lines.append(line.strip())
                    elif line.strip():
                        q1_lines.append(line.strip())
            
                q1_elem = copy.deepcopy(templates['question1'])
                slide.shapes._spTree.append(q1_elem)
                shape = slide.shapes[-1]
                if q1_lines:
                    replace_text_preserve_format(shape, q1_lines, font_color=RGBColor(255, 255, 255), layout_name=layout.name)
                    # If single line, remove bullet formatting from paragraphs
                    if len(q1_lines) == 1:
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        for p_elem in shape.text_frame._txBody.findall('.//a:p', namespaces=ns):
                            pPr = p_elem.find('a:pPr', namespaces=ns)
                            if pPr is not None:
                                for tag in ('a:buFont', 'a:buChar', 'a:buAutoNum'):
                                    el = pPr.find(tag, namespaces=ns)
                                    if el is not None:
                                        pPr.remove(el)
                                pPr.set('marL', '0')
                                pPr.set('indent', '0')
                    print(f"Updated {layout.name} question1 -> {len(q1_lines)} lines")
                else:
                    replace_text_preserve_format(shape, "question1")  # Keep original if missing
                    print(f"Inserted {layout.name} with default text")
                
                # Use centralized metadata injection for topic/subtopic
                apply_metadata_to_slide(slide, {})

        elif layout.name == 'LAYOUT_homework_question_page':
            templates = sst_content_templates.get(layout.name, {})
            # Inject all static elements
            for static_elem in templates.get('static_elements', []):
                elem_copy = remove_locks(copy.deepcopy(static_elem))
                copy_image_rels(elem_copy, layout.part, slide.part)
                slide.shapes._spTree.append(elem_copy)
            
            # Inject dynamic question text
            if templates.get('text'):
                question_text = "\n".join([get_text(entry).strip() for entry in section['content'] if get_text(entry).strip()])
                for text_elem_xml in templates['text']:
                    text_elem = copy.deepcopy(text_elem_xml)
                    slide.shapes._spTree.append(text_elem)
                    shape = slide.shapes[-1]
                    if question_text:
                        replace_text_preserve_format(shape, question_text, layout_name=layout.name)
            
            apply_metadata_to_slide(slide, {})
            print(f"Inserted {layout.name} with dynamic question text")

        elif layout.name == 'LAYOUT_sst_homework_page':
            templates = sst_content_templates.get(layout.name, {})
            # Inject all static elements (groups, text boxes, etc) from layout
            for static_elem in templates.get('static_elements', []):
                elem_copy = remove_locks(copy.deepcopy(static_elem))
                copy_image_rels(elem_copy, layout.part, slide.part)
                slide.shapes._spTree.append(elem_copy)
            
            # Use centralized metadata injection for topic/subtopic
            apply_metadata_to_slide(slide, {})
            print(f"Inserted {layout.name} with all background elements")

        elif layout.name == 'LAYOUT_sst_activity_static_page':
            templates = sst_content_templates.get(layout.name, {})
            # Inject all static elements from layout
            for static_elem in templates.get('static_elements', []):
                elem_copy = remove_locks(copy.deepcopy(static_elem))
                copy_image_rels(elem_copy, layout.part, slide.part)
                slide.shapes._spTree.append(elem_copy)
            
            # Use centralized metadata injection for topic/subtopic
            apply_metadata_to_slide(slide, {})
            print(f"Inserted {layout.name} as static slide")

        # After processing the section and adding elements, check if we need to overlay SYR
        if section.get('has_syr'):
            layout_syr = get_layout('syr')
            if layout_syr:
                syr_templates = sst_content_templates.get('LAYOUT_syr', {})
                for static_elem in syr_templates.get('static_elements', []):
                    elem_copy = remove_locks(copy.deepcopy(static_elem))
                    copy_image_rels(elem_copy, layout_syr.part, slide.part)
                    slide.shapes._spTree.append(elem_copy)
                print(f"Overlaid LAYOUT_syr on current slide.")

        # Check if we need to overlay Ask Question
        ask_q_text = section.get('ask_question_text')
        if ask_q_text:
            layout_ask = get_layout('ask_question')
            if layout_ask:
                ask_templates = sst_content_templates.get('LAYOUT_ask_question', {})
                for static_elem in ask_templates.get('static_elements', []):
                    elem_copy = remove_locks(copy.deepcopy(static_elem))
                    copy_image_rels(elem_copy, layout_ask.part, slide.part)
                    slide.shapes._spTree.append(elem_copy)
                
                # After appending all Ask Question elements, find the group and update its text
                for shp in slide.shapes:
                    if shp.shape_type == 6:  # GROUP
                        for sub in shp.shapes:
                            if getattr(sub, 'has_text_frame', False):
                                if 'write question here' in sub.text.lower():
                                    replace_text_preserve_format(sub, ask_q_text)
                                    # Add 10% padding to left and right inside the text frame
                                    margin = int(sub.width * 0.10)
                                    sub.text_frame.margin_left = margin
                                    sub.text_frame.margin_right = margin
                                    # Center if short (one line), left align if long (multi-line)
                                    align = PP_ALIGN.CENTER if len(ask_q_text) <= 45 else PP_ALIGN.LEFT
                                    for p in sub.text_frame.paragraphs:
                                        p.alignment = align
                print(f"Overlaid LAYOUT_ask_question with text: '{ask_q_text}'")
        
        # Finally, inject logo elements on top of everything else
        inject_logo(slide)

    report_progress(96, "Saving presentation...")
    prs.save(output_path)
    report_progress(100, "Done!")

if __name__ == "__main__":
    generate_ppt("content.docx", "template.pptx", "Generated_Presentation.pptx")
    print("DONE -- PPT Generated", flush=True)
