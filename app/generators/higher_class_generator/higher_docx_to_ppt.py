"""
Higher Class PPT Generator
---------------------------
Generates PowerPoint presentations from a DOCX file using the
higher-class template (template_higher_class.pptx).

Supported DOCX section markers (bracketed tags):
  [science_title_page]       – Title slide with class, chapter, topic metadata
  [science_lo_page]          – Learning Objectives slide
  [science_transition_page]  – Transition / "Let's discuss" slide with topic
  [science_content_page]     – Generic content slide
  [science_discussion_page]  – Discussion slide
  [science_quiztime_page]    – Quiz-time slide (auto-selects page1 vs page2)
  [finalquiz_page]           – Final Quiz intro + Q/A pairs
  [syr]                      – "Show Your Response" overlay marker
"""

import re
import os
import copy
import io
from docx import Document
from docx.document import Document as _Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml import parse_xml
from PIL import Image


def generate_ppt(docx_path, template_path, output_path, progress_callback=None):
    def report_progress(percentage, status):
        if progress_callback:
            progress_callback(percentage, status)

    report_progress(2, "Initializing...")
    prs = Presentation(template_path)
    doc = Document(docx_path)

    def clean(text):
        return re.sub(r'[^a-z0-9]', '', text.strip().lower())

    global_metadata = {
        'class': '',
        'subject': '',
        'chapter number': '',
        'chapter name': '',
        'lesson': '',
        'lesson number': '',
        'topic': '',
        'subtopic': ''
    }

    def get_text(entry):
        """Extract plain text from a content entry (either a string or a (text_parts, ilvl) tuple)."""
        data = entry[0] if isinstance(entry, tuple) else entry
        if isinstance(data, list):
            return "".join(p['value'] if p['type'] == 'text' else '' for p in data)
        return str(data)

    def get_content_parts(entry):
        """Extract the list of content parts from a content entry."""
        data = entry[0] if isinstance(entry, tuple) else entry
        if isinstance(data, list):
            return data
        return [{'type': 'text', 'value': str(data)}]

    def copy_image_rels(element, source_part, target_part):
        """Copy image relationships referenced by ANY element (blip, svgBlip, etc.) from source to target part."""
        r_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
        
        # Iterate over all elements including the root itself
        for el in [element] + list(element.iter()):
            for attr_name in (f'{r_ns}embed', f'{r_ns}link'):
                old_rid = el.get(attr_name)
                if old_rid and old_rid in source_part.rels:
                    rel = source_part.rels[old_rid]
                    try:
                        # relate_to handles duplicates and returns the existing rId if the same part is already linked
                        new_rid = target_part.relate_to(rel.target_part, rel.reltype)
                        el.set(attr_name, new_rid)
                    except Exception as e:
                        print(f"Warning: Could not copy relationship {old_rid}: {e}")

    def remove_locks(element):
        """Remove a:spLocks and a:grpSpLocks elements from a shape XML to make it fully editable."""
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        for lock in element.findall('.//a:spLocks', namespaces=ns):
            lock.getparent().remove(lock)
        for lock in element.findall('.//a:grpSpLocks', namespaces=ns):
            lock.getparent().remove(lock)
        return element

    def apply_locks(element):
        """Add protection locks to all shapes within an element to prevent interaction."""
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
              'p': 'http://schemas.openxmlformats.org/officeDocument/2006/main',
              'p_main': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        targets = [
            ('.//p_main:nvSpPr', 'p_main:cNvSpPr', 'a:spLocks'),
            ('.//p_main:nvPicPr', 'p_main:cNvPicPr', 'a:picLocks'),
            ('.//p_main:nvGrpSpPr', 'p_main:cNvGrpSpPr', 'a:grpSpLocks'),
            ('.//p_main:nvCxnSpPr', 'p_main:cNvCxnSpPr', 'a:cxnSpLocks')
        ]
        for search_xpath, cNvPr_tag, lock_tag in targets:
            candidates = []
            candidates.extend(element.findall(search_xpath, namespaces=ns))
            root_tag_local = search_xpath.split(':')[-1]
            if f"{{{ns['p_main']}}}{root_tag_local}" == element.tag:
                candidates.append(element)
            for nvPr in candidates:
                cNvPr = nvPr.find(cNvPr_tag, namespaces=ns)
                if cNvPr is not None:
                    locks = cNvPr.find(lock_tag, namespaces=ns)
                    if locks is None:
                        locks = parse_xml(f'<a:{lock_tag.split(":")[1]} xmlns:a="{ns["a"]}" noGrp="1" noSelect="1" noRot="1" noChangeAspect="1" noMove="1" noResize="1" noEditPoints="1" noAdjustHandles="1" noChangeArrowheads="1" noChangeShapeType="1" noTextEdit="1"/>')
                        cNvPr.append(locks)
                    else:
                        for attr in ("noGrp", "noSelect", "noRot", "noChangeAspect", "noMove", "noResize", "noEditPoints", "noAdjustHandles", "noChangeArrowheads", "noChangeShapeType", "noTextEdit"):
                            locks.set(attr, "1")
        return element

    # ── Layout lookup ───────────────────────────────────────────────────
    def get_layout(name):
        """Find a slide layout by name. Supports several alias patterns."""
        # Exact match with LAYOUT_ prefix
        for layout in prs.slide_layouts:
            if layout.name == f"LAYOUT_{name}":
                return layout

        # Case-insensitive, underscore-insensitive match
        target = f"layout_{name}".replace("_", "").lower()
        for layout in prs.slide_layouts:
            if layout.name.replace("_", "").replace(" ", "").lower() == target:
                return layout

        # Science-specific aliases
        cleaned = name.replace("_", "").lower()

        if cleaned in ('sciencetitlepage', 'sciencepagetitle'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_science_title_page':
                    return layout

        if cleaned in ('sciencelopage', 'learningobjective', 'sciencelo'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_science_lo_page':
                    return layout

        if cleaned in ('sciencetransitionpage', 'transitionpage', 'transition'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_science_transition_page':
                    return layout

        if cleaned in ('sciencecontentpage', 'contentpage', 'content'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_science_content_page':
                    return layout

        if cleaned in ('sciencediscussionpage', 'discussionpage', 'discussion'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_science_discussion_page':
                    return layout

        if cleaned in ('sciencequiztimepage', 'quiztimepage', 'quiztime'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_science_quiztime_q_page1':
                    return layout

        if cleaned in ('finalquizpage', 'finalquiz', 'quiz'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_finalquiz_page':
                    return layout

        if cleaned == 'syr':
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_syr':
                    return layout

        return None

    # ── Template dictionaries ───────────────────────────────────────────
    # Store extracted shapes from layouts for later injection
    content_templates = {
        'LAYOUT_science_content_page': {'topic': None, 'subtopic': None, 'text': [], 'static_elements': []},
        'LAYOUT_science_discussion_page': {'question_text': None, 'static_elements': []},
        'LAYOUT_science_quiztime_q_page1': {'question': None, 'options': [], 'option_labels': [], 'static_elements': []},
        'LAYOUT_science_quiztime_a_page1': {'question': None, 'options': [], 'option_labels': [], 'static_elements': []},
        'LAYOUT_science_quiztime_q_page2': {'question': None, 'option_groups': [], 'static_elements': []},
        'LAYOUT_science_quiztime_a_page2': {'question': None, 'option_groups': [], 'static_elements': []},
        'LAYOUT_syr': {'static_elements': []},
        'LAYOUT_finalquiz_page': {'static_elements': []},
    }

    # Store LO page elements: group template (oval+rect), title, connector
    lo_page_data = {
        'group_xml': None,         # The group with oval(number) + rounded-rect(text)
        'title_xml': None,         # "In this lesson, we will explore"
        'connector_xml': None,     # Straight Connector decoration
        'layout_part': None,       # For image rels
    }

    # Store Discussion page elements: header (Group 6), question template (Group 28), subtitle (TextBox 31)
    discussion_page_data = {
        'header_xml': None,        # Group 6: "Discussion Time" + line
        'question_grp_xml': None,  # Group 28: Rounded rect + icon
        'subtitle_xml': None,      # TextBox 31: "Discuss with your partner:"
        'layout_part': None,       # For image rels
    }

    # Capture logo elements (if LAYOUT_logo exists in template)
    logo_elements = []
    logo_source_part = None

    def inject_logo(target_slide):
        """Inject captured logo elements onto a target slide and lock them."""
        for logo_xml in logo_elements:
            new_logo_elem = apply_locks(copy.deepcopy(logo_xml))
            target_slide.shapes._spTree.append(new_logo_elem)
            if logo_source_part:
                copy_image_rels(new_logo_elem, logo_source_part, target_slide.part)

    # ── Pre-process layouts: extract template shapes ────────────────────
    for layout in prs.slide_layouts:
        # LO page: extract group (oval + rounded rect), title rect, and connector
        if layout.name == 'LAYOUT_science_lo_page':
            lo_page_data['layout_part'] = layout.part
            for shape in list(layout.shapes):
                if shape.shape_type == 6:  # Group
                    lo_page_data['group_xml'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                elif shape.shape_type == 9:  # Line/Connector
                    lo_page_data['connector_xml'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                elif getattr(shape, 'has_text_frame', False) and shape.text.strip():
                    lo_page_data['title_xml'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)

        # Logo layout
        if layout.name == "LAYOUT_logo":
            logo_source_part = layout.part
            for shape in list(layout.shapes):
                logo_elements.append(copy.deepcopy(shape.element))
                shape.element.getparent().remove(shape.element)

        # SYR / finalquiz_page static elements
        if layout.name in ('LAYOUT_syr', 'LAYOUT_finalquiz_page'):
            templates = content_templates.get(layout.name, {})
            for shape in list(layout.shapes):
                if getattr(shape, 'is_placeholder', False) and shape.placeholder_format.type == 18:
                    continue
                templates['static_elements'].append(copy.deepcopy(shape.element))
                shape.element.getparent().remove(shape.element)

        # Discussion page: capture specific elements
        if layout.name == 'LAYOUT_science_discussion_page':
            discussion_page_data['layout_part'] = layout.part
            for shape in list(layout.shapes):
                if shape.name == "Group 6":
                    discussion_page_data['header_xml'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                elif shape.name == "Group 28":
                    discussion_page_data['question_grp_xml'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                elif shape.name == "TextBox 31":
                    discussion_page_data['subtitle_xml'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                elif shape.shape_type == 6: # Catch-all for any other groups
                    discussion_page_data['header_xml'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)

        # Quiztime pages: capture question, options, and labels
        if layout.name in ('LAYOUT_science_quiztime_q_page1', 'LAYOUT_science_quiztime_a_page1'):
            templates = content_templates[layout.name]
            for shape in list(layout.shapes):
                if getattr(shape, 'has_text_frame', False):
                    txt = shape.text.strip()
                    # Long question text (not a single-letter label)
                    if len(txt) > 2 and shape.shape_type == 17:  # TextBox with question
                        templates['question'] = copy.deepcopy(shape.element)
                        shape.element.getparent().remove(shape.element)
                    elif len(txt) == 1 and txt in 'ABCD':  # Option labels
                        templates['option_labels'].append(copy.deepcopy(shape.element))
                        shape.element.getparent().remove(shape.element)
                    elif shape.shape_type == 1:  # AUTO_SHAPE (rounded rectangle options)
                        templates['options'].append(copy.deepcopy(shape.element))
                        shape.element.getparent().remove(shape.element)
                    else:
                        templates['static_elements'].append(copy.deepcopy(shape.element))
                        shape.element.getparent().remove(shape.element)
                elif shape.shape_type == 6:  # Group (decorative)
                    templates['static_elements'].append(copy.deepcopy(shape.element))
                    shape.element.getparent().remove(shape.element)
                else:
                    templates['static_elements'].append(copy.deepcopy(shape.element))
                    shape.element.getparent().remove(shape.element)

        if layout.name in ('LAYOUT_science_quiztime_q_page2', 'LAYOUT_science_quiztime_a_page2'):
            templates = content_templates[layout.name]
            for shape in list(layout.shapes):
                if getattr(shape, 'has_text_frame', False) and len(shape.text.strip()) > 2:
                    templates['question'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                elif shape.shape_type == 6:  # Groups (option groups)
                    templates['option_groups'].append(copy.deepcopy(shape.element))
                    shape.element.getparent().remove(shape.element)
                else:
                    templates['static_elements'].append(copy.deepcopy(shape.element))
                    shape.element.getparent().remove(shape.element)

    # ── Parse DOCX ──────────────────────────────────────────────────────
    def iter_block_items(parent):
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            raise ValueError("iter_block_items: parent must be Document or _Cell")
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    sections = []
    current_section = None
    w_ns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'

    for item in iter_block_items(doc):
        lines_to_process = []
        images_to_add = []

        if isinstance(item, Paragraph):
            ilvl = 0
            is_list = False
            numPr = item._element.find(f'.//{{{w_ns}}}numPr')
            if numPr is not None:
                is_list = True
                ilvl_elem = numPr.find(f'{{{w_ns}}}ilvl')
                if ilvl_elem is not None:
                    ilvl = int(ilvl_elem.get(f'{{{w_ns}}}val', '0'))
            else:
                if item.style and item.style.name and ('List' in item.style.name or 'Bullet' in item.style.name):
                    is_list = True
                    import re as _re
                    m = _re.search(r'\d+', item.style.name)
                    if m:
                        ilvl = max(0, int(m.group(0)) - 1)

            m_ns = {'m': 'http://schemas.openxmlformats.org/officeDocument/2006/math'}
            parts = []

            for child in item._element.iterchildren():
                if child.tag.endswith('}oMath'):
                    from lxml import etree
                    math_xml = etree.tostring(child, encoding='unicode')
                    parts.append({'type': 'math', 'value': math_xml})
                elif child.tag.endswith('}oMathPara'):
                    for subchild in child.iterchildren():
                        if subchild.tag.endswith('}oMath'):
                            from lxml import etree
                            math_xml = etree.tostring(subchild, encoding='unicode')
                            parts.append({'type': 'math', 'value': math_xml})
                elif child.tag.endswith('}r'):
                    t_elem = child.find(f'.//{{{w_ns}}}t')
                    if t_elem is not None and t_elem.text:
                        rPr = child.find(f'.//{{{w_ns}}}rPr')
                        is_bold = False
                        is_italic = False
                        if rPr is not None:
                            b_elem = rPr.find(f'.//{{{w_ns}}}b')
                            i_elem = rPr.find(f'.//{{{w_ns}}}i')
                            if b_elem is not None:
                                val = b_elem.get(f'{{{w_ns}}}val', '1')
                                is_bold = val not in ('0', 'false', 'False')
                            if i_elem is not None:
                                val = i_elem.get(f'{{{w_ns}}}val', '1')
                                is_italic = val not in ('0', 'false', 'False')
                        parts.append({'type': 'text', 'value': t_elem.text, 'bold': is_bold, 'italic': is_italic})

            if not parts:
                lines_to_process.append((item.text.strip(), ilvl, is_list))
            else:
                lines_to_process.append((parts, ilvl, is_list))

            for run in item.runs:
                for drawing in run._element.findall('.//w:drawing', namespaces=run._element.nsmap):
                    for blip in drawing.findall('.//a:blip', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                        embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if embed:
                            blob = doc.part.related_parts[embed].blob
                            images_to_add.append(blob)
        elif isinstance(item, Table):
            for row in item.rows:
                cells = [c.text.strip() for c in row.cells]
                if len(cells) >= 2:
                    key = cells[0].rstrip(':').strip()
                    val = cells[1].lstrip(':').strip()
                    lines_to_process.append((f"{key}: {val}", 0))
                elif len(cells) == 1:
                    lines_to_process.append((cells[0], 0))

        for entry in lines_to_process:
            content_data = entry[0] if isinstance(entry, tuple) else entry
            ilvl = entry[1] if isinstance(entry, tuple) else 0

            if not content_data and not images_to_add:
                continue

            if isinstance(content_data, list):
                search_text = "".join(p['value'] if p['type'] == 'text' else '' for p in content_data)
            else:
                search_text = str(content_data)

            match = re.search(r'^\[\s*([^\]]+?)\s*\]', search_text)
            if match:
                lower_text = search_text.lower()
                has_syr = '[add_syr]' in lower_text or '[syr]' in lower_text

                name = match.group(1).strip()

                current_section = {
                    'name': name,
                    'content': [],
                    'images': [],
                    'has_syr': has_syr,
                }
                sections.append(current_section)
            elif current_section is not None:
                if content_data:
                    lower_text = search_text.lower()
                    if '[add_syr]' in lower_text or '[syr]' in lower_text:
                        current_section['has_syr'] = True
                        if isinstance(content_data, str):
                            content_data = re.sub(r'\[add_syr\]|\[syr\]', '', content_data, flags=re.IGNORECASE).strip()

                    if content_data:
                        current_section['content'].append((content_data, ilvl))
                if images_to_add:
                    current_section['images'].extend(images_to_add)
                    images_to_add = []

    report_progress(10, "Document parsed. Generating slides...")

    # ── Text helpers ────────────────────────────────────────────────────
    def replace_text_preserve_format(shape, new_text, center=False, font_color=None, layout_name=None):
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        if not tf.paragraphs or not tf.paragraphs[0].runs:
            if isinstance(new_text, list):
                shape.text = "\n".join([str(t[0] if isinstance(t, tuple) else t) for t in new_text])
            else:
                shape.text = str(new_text)
            if center:
                for p in tf.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
            return

        p0 = tf.paragraphs[0]
        p0_xml = copy.deepcopy(p0._p)

        font = p0.runs[0].font
        fname = font.name
        fsize = font.size
        fbold = font.bold
        fitalic = font.italic
        funderline = font.underline

        try:
            color_type = getattr(font.color, 'type', None)
            color_rgb = getattr(font.color, 'rgb', None) if color_type == 1 else None
            color_theme = getattr(font.color, 'theme_color', None) if color_type == 2 else None
        except:
            color_type, color_rgb, color_theme = None, None, None

        tf.clear()

        texts = new_text if isinstance(new_text, list) else [new_text]

        for i, para_data in enumerate(texts):
            if i == 0:
                p = tf.paragraphs[0]
                new_p_xml = copy.deepcopy(p0_xml)
                p._p.getparent().replace(p._p, new_p_xml)
                p = tf.paragraphs[0]
            else:
                new_p_xml = copy.deepcopy(p0_xml)
                tf._txBody.append(new_p_xml)
                p = tf.paragraphs[i]
                p.space_before = Pt(12)

            level = 0
            is_list = False
            if isinstance(para_data, tuple):
                level = para_data[1]
                if len(para_data) > 2:
                    is_list = para_data[2]
                para_data = para_data[0]
            p.level = level

            if is_list:
                a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                pPr = p._p.get_or_add_pPr()
                for child in list(pPr):
                    if child.tag.endswith('}buNone') or child.tag.endswith('}buAutoNum') or child.tag.endswith('}buChar') or child.tag.endswith('}buFont'):
                        pPr.remove(child)
                from lxml.etree import fromstring as parse_xml
                char_val = '•' if level == 0 else ('-' if level == 1 else '▪')
                buChar_xml = f'<a:buChar xmlns:a="{a_ns}" char="{char_val}"/>'
                pPr.insert(0, parse_xml(buChar_xml))

            for r_elem in p._p.findall('.//a:r', namespaces=p._p.nsmap):
                p._p.remove(r_elem)
            for fld_elem in p._p.findall('.//a:fld', namespaces=p._p.nsmap):
                p._p.remove(fld_elem)

            if center:
                p.alignment = PP_ALIGN.CENTER

            parts = para_data if isinstance(para_data, list) else [{'type': 'text', 'value': str(para_data)}]

            # Layouts that should have bold text colored yellow (#FFC000)
            yellow_bold_layouts = [
                'LAYOUT_science_lo_page',
            ]

            for part in parts:
                if part['type'] == 'text':
                    new_run = p.add_run()
                    new_run.text = part['value']

                    if fname is not None:
                        new_run.font.name = fname
                    if fsize is not None:
                        new_run.font.size = fsize

                    part_bold = part.get('bold', False)
                    if part_bold or fbold is not None:
                        new_run.font.bold = part_bold or fbold

                    part_italic = part.get('italic', False)
                    if part_italic or fitalic is not None:
                        new_run.font.italic = part_italic or fitalic

                    if funderline is not None:
                        new_run.font.underline = funderline

                    if part_bold and layout_name in yellow_bold_layouts:
                        new_run.font.color.rgb = RGBColor(255, 192, 0)
                    elif font_color is not None:
                        new_run.font.color.rgb = font_color
                    else:
                        if color_rgb is not None:
                            new_run.font.color.rgb = color_rgb
                        elif color_theme is not None:
                            new_run.font.color.theme_color = color_theme
                elif part['type'] == 'math':
                    try:
                        from lxml.etree import fromstring as lxml_parse
                        math_elem = lxml_parse(part['value'])
                        p._p.append(math_elem)
                    except Exception as e:
                        print(f"Error injecting math XML: {e}")

    def apply_metadata_to_slide(slide, slide_data):
        """Recursively process all shapes in a slide and update metadata placeholders."""
        merged_data = {**global_metadata, **slide_data}

        def process_shape_list(shapes, parent_group=None):
            for shape in shapes:
                if shape.shape_type == 6:
                    process_shape_list(shape.shapes, parent_group=shape)
                elif getattr(shape, 'has_text_frame', False):
                    cleaned_txt = clean(shape.text)
                    if not cleaned_txt:
                        continue

                    if 'subtopic' in cleaned_txt:
                        txt_low = shape.text.lower()
                        if 'topic' in txt_low.replace('subtopic', ''):
                            topic_val = merged_data.get('topic', '')
                            subtopic_val = merged_data.get('subtopic', '')
                            if topic_val and subtopic_val:
                                val = f"{topic_val} - {subtopic_val}"
                            else:
                                val = topic_val or subtopic_val
                        else:
                            val = merged_data.get('subtopic')

                        if val:
                            replace_text_preserve_format(shape, val, center=True, layout_name=slide.slide_layout.name)
                            tf = shape.text_frame
                            tf.word_wrap = False
                            text_width = len(str(val).strip()) * 250000
                            padding = 800000
                            required_width = text_width + padding
                            if shape.width > 0 and required_width < shape.width:
                                required_width = shape.width
                            max_w = int(prs.slide_width * 0.95)
                            if required_width > max_w:
                                required_width = max_w
                                tf.word_wrap = True
                            tf.margin_left = int(padding / 2.5)
                            tf.margin_right = int(padding / 2.5)
                            tf.margin_top = Pt(4)
                            tf.margin_bottom = Pt(4)
                            if shape.width > 0:
                                scale = required_width / shape.width
                                if parent_group:
                                    g_left = parent_group.left
                                    for child in parent_group.shapes:
                                        child.width = int(child.width * scale)
                                        rel_left = child.left - g_left
                                        child.left = g_left + int(rel_left * scale)
                                    parent_group.width = int(parent_group.width * scale)
                                else:
                                    shape.width = required_width
                            else:
                                shape.width = required_width

                    elif 'topic' in cleaned_txt and merged_data.get('topic'):
                        replace_text_preserve_format(shape, merged_data['topic'], center=True, layout_name=slide.slide_layout.name)

                    elif cleaned_txt == 'class' and merged_data.get('class'):
                        replace_text_preserve_format(shape, merged_data['class'], center=True, layout_name=slide.slide_layout.name)
                    elif cleaned_txt == 'subject' and merged_data.get('subject'):
                        replace_text_preserve_format(shape, merged_data['subject'], center=True, layout_name=slide.slide_layout.name)
                    elif cleaned_txt == 'chapternumber' and merged_data.get('chapter number'):
                        replace_text_preserve_format(shape, merged_data['chapter number'], center=True, layout_name=slide.slide_layout.name)
                    elif cleaned_txt == 'chaptername' and merged_data.get('chapter name'):
                        replace_text_preserve_format(shape, merged_data['chapter name'], center=True, layout_name=slide.slide_layout.name)
                    elif cleaned_txt in ('lesson', 'lessonnumber') and (merged_data.get('lesson') or merged_data.get('lesson number')):
                        val = merged_data.get('lesson') or merged_data.get('lesson number')
                        replace_text_preserve_format(shape, val, center=True, layout_name=slide.slide_layout.name)

        process_shape_list(slide.shapes)

    # ── Generate slides ─────────────────────────────────────────────────
    total_sections = len(sections)
    for i, section in enumerate(sections):
        progress = 10 + int((i / total_sections) * 85)
        report_progress(progress, f"Generating slide {i+1} of {total_sections}: {section['name']}...")

        sname = section['name'].strip().lower()

        # ── Dynamic layout selection for quiz ───────────────────────
        if sname in ('science_quiztime_page', 'quiztime_page', 'quiztime'):
            quiz_data = {'question': '', 'options': []}
            for entry in section['content']:
                line = get_text(entry)
                line_low = line.lower()
                if 'question:' in line_low:
                    quiz_data['question'] = line.split(':', 1)[1].strip()
                elif 'options:' in line_low:
                    opts = line.split(':', 1)[1].strip()
                    quiz_data['options'] = [o.strip() for o in opts.split(',')]
                elif line.strip():
                    if not quiz_data['question']:
                        quiz_data['question'] = line.strip()
                    else:
                        quiz_data['options'].append(line.strip())

            # page1 layout has separate A/B/C/D option boxes (short text)
            # page2 layout has group-based options (long text)
            use_long = any(len(opt) > 25 for opt in quiz_data['options'])
            if use_long:
                q_layout_name = 'LAYOUT_science_quiztime_q_page2'
                a_layout_name = 'LAYOUT_science_quiztime_a_page2'
            else:
                q_layout_name = 'LAYOUT_science_quiztime_q_page1'
                a_layout_name = 'LAYOUT_science_quiztime_a_page1'

            # Find layouts
            q_layout = None
            a_layout = None
            for lay in prs.slide_layouts:
                if lay.name == q_layout_name:
                    q_layout = lay
                if lay.name == a_layout_name:
                    a_layout = lay

            if q_layout:
                q_slide = prs.slides.add_slide(q_layout)
                q_templates = content_templates.get(q_layout_name, {})

                # Inject static elements
                for static_elem in q_templates.get('static_elements', []):
                    elem_copy = remove_locks(copy.deepcopy(static_elem))
                    copy_image_rels(elem_copy, q_layout.part, q_slide.part)
                    q_slide.shapes._spTree.append(elem_copy)

                # Inject question
                if q_templates.get('question') is not None:
                    q_elem = copy.deepcopy(q_templates['question'])
                    q_slide.shapes._spTree.append(q_elem)
                    replace_text_preserve_format(q_slide.shapes[-1], quiz_data['question'])

                if not use_long:
                    # page1: inject option labels and option boxes
                    for label_elem in q_templates.get('option_labels', []):
                        q_slide.shapes._spTree.append(copy.deepcopy(label_elem))
                    for idx, opt_elem in enumerate(q_templates.get('options', [])):
                        elem = copy.deepcopy(opt_elem)
                        q_slide.shapes._spTree.append(elem)
                        if idx < len(quiz_data['options']):
                            replace_text_preserve_format(q_slide.shapes[-1], quiz_data['options'][idx], center=True)
                        else:
                            replace_text_preserve_format(q_slide.shapes[-1], '', center=True)
                else:
                    # page2: inject option groups
                    for idx, grp_elem in enumerate(q_templates.get('option_groups', [])):
                        elem = copy.deepcopy(grp_elem)
                        q_slide.shapes._spTree.append(elem)
                        grp_shape = q_slide.shapes[-1]
                        if idx < len(quiz_data['options']):
                            for sub in grp_shape.shapes:
                                if getattr(sub, 'has_text_frame', False) and sub.text.strip():
                                    replace_text_preserve_format(sub, quiz_data['options'][idx], center=True)
                                    break

                inject_logo(q_slide)
                print(f"Generated quiz question slide ({q_layout_name})")

            if a_layout:
                a_slide = prs.slides.add_slide(a_layout)
                a_templates = content_templates.get(a_layout_name, {})

                # Inject static elements
                for static_elem in a_templates.get('static_elements', []):
                    elem_copy = remove_locks(copy.deepcopy(static_elem))
                    copy_image_rels(elem_copy, a_layout.part, a_slide.part)
                    a_slide.shapes._spTree.append(elem_copy)

                # Inject question
                if a_templates.get('question') is not None:
                    a_elem = copy.deepcopy(a_templates['question'])
                    a_slide.shapes._spTree.append(a_elem)
                    replace_text_preserve_format(a_slide.shapes[-1], quiz_data['question'])

                if not use_long:
                    for label_elem in a_templates.get('option_labels', []):
                        a_slide.shapes._spTree.append(copy.deepcopy(label_elem))
                    for idx, opt_elem in enumerate(a_templates.get('options', [])):
                        elem = copy.deepcopy(opt_elem)
                        a_slide.shapes._spTree.append(elem)
                        if idx < len(quiz_data['options']):
                            replace_text_preserve_format(a_slide.shapes[-1], quiz_data['options'][idx], center=True)
                        else:
                            replace_text_preserve_format(a_slide.shapes[-1], '', center=True)
                else:
                    for idx, grp_elem in enumerate(a_templates.get('option_groups', [])):
                        elem = copy.deepcopy(grp_elem)
                        a_slide.shapes._spTree.append(elem)
                        grp_shape = a_slide.shapes[-1]
                        if idx < len(quiz_data['options']):
                            for sub in grp_shape.shapes:
                                if getattr(sub, 'has_text_frame', False) and sub.text.strip():
                                    replace_text_preserve_format(sub, quiz_data['options'][idx], center=True)
                                    break

                inject_logo(a_slide)
                print(f"Generated quiz answer slide ({a_layout_name})")

            # SYR overlay on quiz slides if requested
            if section.get('has_syr'):
                layout_syr = get_layout('syr')
                if layout_syr:
                    syr_templates = content_templates.get('LAYOUT_syr', {})
                    for target_slide in [q_slide, a_slide]:
                        if target_slide:
                            for static_elem in syr_templates.get('static_elements', []):
                                elem_copy = remove_locks(copy.deepcopy(static_elem))
                                copy_image_rels(elem_copy, layout_syr.part, target_slide.part)
                                target_slide.shapes._spTree.append(elem_copy)
                    print("Overlaid LAYOUT_syr on quiz slides.")

            continue  # Skip the rest of the loop for quiztime

        # ── Standard layout selection ───────────────────────────────
        layout = get_layout(section['name'])
        if not layout:
            print(f"Skipping section [{section['name']}], layout not found.")
            continue

        slide = prs.slides.add_slide(layout)

        # ── science_title_page ──────────────────────────────────────
        if section['name'].replace("_", "").lower() in ('sciencetitlepage', 'sciencepagetitle'):
            data = {}
            for entry in section['content']:
                line = get_text(entry)
                if ":" in line:
                    parts = line.split(":", 1)
                    key = clean(parts[0])
                    val = parts[1].strip()
                    data[key.upper()] = val
                    for g_key in global_metadata:
                        if clean(g_key) == key:
                            global_metadata[g_key] = val
                            break

            # Handle combined 'class_&_subject' key (cleaned to 'CLASSSUBJECT')
            # The DOCX may use "class_&_subject: Class 11 Physics" as a single key
            combined_key = None
            for k in list(data.keys()):
                if 'CLASSSUBJECT' in k or 'CLASSANDSUBJECT' in k:
                    combined_key = k
                    break
            if combined_key:
                combined_val = data[combined_key]
                # Store full combined value as 'class' for the template placeholder
                data['CLASS'] = combined_val
                global_metadata['class'] = combined_val
                # Also try to extract subject: look for known subject names at the end
                import re as _re
                subj_match = _re.search(r'\b(Physics|Chemistry|Biology|Mathematics|Science|Computer Science|English|Hindi|Odia)\s*$', combined_val, flags=_re.IGNORECASE)
                if subj_match:
                    data['SUBJECT'] = subj_match.group(1).strip()
                    global_metadata['subject'] = subj_match.group(1).strip()
                    # Extract just the class part (everything before the subject)
                    class_part = combined_val[:subj_match.start()].strip()
                    if class_part:
                        data['CLASS'] = class_part
                        global_metadata['class'] = class_part

            idx_mapping = {}
            for shape in layout.shapes:
                if shape.is_placeholder and shape.has_text_frame:
                    idx_mapping[clean(shape.text)] = shape.placeholder_format.idx

            mapping = {
                "CLASS": "class & subject",
                "CLASSSUBJECT": "class & subject",
                "SUBJECT": "subject",
                "CHAPTER_NUMBER": "chapter number",
                "CHAPTERNUMBER": "chapter number",
                "CHAPTER_NAME": "chapter name",
                "CHAPTERNAME": "chapter name",
                "LESSON": "lesson",
                "LESSON_NUMBER": "lesson number",
                "LESSONNUMBER": "lesson number",
                "TOPIC": "topic"
            }

            for key, template_word in mapping.items():
                cleaned_word = clean(template_word)
                val = data.get(key) or data.get(key.replace("_", "")) or global_metadata.get(template_word)
                if val:
                    if cleaned_word in idx_mapping:
                        idx = idx_mapping[cleaned_word]
                        found_on_slide = False
                        for shape in slide.shapes:
                            if shape.is_placeholder and shape.placeholder_format.idx == idx:
                                shape.text = val
                                print(f"Updated TitlePage: {template_word} -> {val}")
                                found_on_slide = True
                                break
                        if not found_on_slide:
                            print(f"[{template_word}] idx {idx} NOT FOUND on slide shapes!")
                    else:
                        # Not an error for 'subject' — template may not have a subject placeholder
                        if template_word != 'subject':
                            print(f"[{template_word}] NOT FOUND in idx_mapping! keys={idx_mapping.keys()}")
                else:
                    if template_word != 'subject':
                        print(f"[{key}] NOT FOUND in docx data or global metadata!")

            apply_metadata_to_slide(slide, data)

        # ── science_lo_page ─────────────────────────────────────────
        elif layout.name == 'LAYOUT_science_lo_page':
            a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

            # 1. Inject static title ("In this lesson, we will explore")
            if lo_page_data['title_xml'] is not None:
                slide.shapes._spTree.append(copy.deepcopy(lo_page_data['title_xml']))

            # 2. Inject connector line
            if lo_page_data['connector_xml'] is not None:
                slide.shapes._spTree.append(copy.deepcopy(lo_page_data['connector_xml']))

            # 3. Parse content sentences from DOCX
            lo_sentences = []
            for entry in section['content']:
                raw = entry[0] if isinstance(entry, tuple) else entry
                if isinstance(raw, list):
                    plain = "".join([p['value'] for p in raw if p['type'] == 'text']).strip()
                else:
                    plain = str(raw).strip()
                # Strip leading bullets/numbers
                plain = re.sub(r'^[\d\.\)\-\*•]+\s*', '', plain).strip()
                if plain:
                    lo_sentences.append(raw)

            # 4. Clone the group template for each sentence
            group_xml = lo_page_data['group_xml']
            if group_xml is not None and lo_sentences:
                # Determine vertical spacing
                # The group's original Y position is ~2089532 EMU (from template)
                # Each cloned group is offset downward
                row_height = int(Inches(1.1))  # spacing between each numbered point

                for idx, sentence_data in enumerate(lo_sentences):
                    new_grp = copy.deepcopy(group_xml)

                    # Copy image rels if needed
                    if lo_page_data['layout_part']:
                        copy_image_rels(new_grp, lo_page_data['layout_part'], slide.part)

                    slide.shapes._spTree.append(new_grp)
                    grp_shape = slide.shapes[-1]

                    # Shift the group down for each subsequent item
                    grp_shape.top = grp_shape.top + (idx * row_height)

                    # Update the sub-shapes within the cloned group
                    for sub in grp_shape.shapes:
                        if not getattr(sub, 'has_text_frame', False):
                            continue
                        sub_text = sub.text.strip().lower()

                        if sub_text in ('1', '2', '3', '4', '5', '6', '7', '8', '9', '10'):
                            # This is the number oval — update the number
                            tf = sub.text_frame
                            for p in tf.paragraphs:
                                for run in p.runs:
                                    run.text = str(idx + 1)

                        elif sub_text == 'text' or len(sub_text) > 1:
                            # This is the content rounded rectangle — replace with sentence
                            tf = sub.text_frame
                            txBody = tf._txBody

                            # Get first paragraph as template, clear its runs
                            p0 = tf.paragraphs[0]
                            for r_elem in list(p0._p.findall(f'{{{a_ns}}}r')):
                                p0._p.remove(r_elem)

                            # Build content parts
                            parts = get_content_parts(sentence_data)
                            for part in parts:
                                if part['type'] == 'text':
                                    run = p0.add_run()
                                    run.text = part['value']
                                    run.font.size = Pt(28)
                                    run.font.name = 'Calibri'
                                    part_bold = part.get('bold', False)
                                    if part_bold:
                                        run.font.bold = True
                                        run.font.color.rgb = RGBColor(255, 192, 0)  # Yellow for bold
                                    else:
                                        run.font.color.rgb = RGBColor(255, 255, 255)  # White
                                elif part['type'] == 'math':
                                    try:
                                        from lxml.etree import fromstring as lxml_parse
                                        math_elem = lxml_parse(part['value'])
                                        p0._p.append(math_elem)
                                    except Exception as e:
                                        print(f"Error injecting math: {e}")

                print(f"Inserted {layout.name} with {len(lo_sentences)} numbered points")
            else:
                if not lo_sentences:
                    print(f"WARNING: No content found for {layout.name}")
                else:
                    print(f"WARNING: No group template found for {layout.name}")

            apply_metadata_to_slide(slide, {})

        # ── science_transition_page ─────────────────────────────────
        elif layout.name == 'LAYOUT_science_transition_page':
            # This layout has a placeholder for "topic" (idx 10) and a
            # "Let's discuss" static text shape (keep constant).
            # The content lines become the topic text.
            topic_text = ''
            for entry in section['content']:
                line = get_text(entry)
                if line.strip():
                    topic_text = line.strip()
                    break

            if topic_text:
                global_metadata['topic'] = topic_text

            # Directly set placeholder idx 10 (topic) on the slide
            if topic_text:
                for shape in slide.shapes:
                    if shape.is_placeholder and shape.placeholder_format.idx == 10:
                        replace_text_preserve_format(shape, topic_text, center=True, layout_name=layout.name)
                        print(f"Set transition topic placeholder (idx 10) -> '{topic_text}'")
                        break

            apply_metadata_to_slide(slide, {'topic': topic_text})
            print(f"Inserted {layout.name} with topic: '{topic_text}'")

        # ── science_content_page ────────────────────────────────────
        elif layout.name == 'LAYOUT_science_content_page':
            # This is an empty layout — content is placed via text boxes
            data = {}
            data_text_list = []
            for entry in section['content']:
                content_obj = entry[0] if isinstance(entry, tuple) else entry
                ilvl = entry[1] if isinstance(entry, tuple) else 0
                is_list = entry[2] if (isinstance(entry, tuple) and len(entry) > 2) else False
                if isinstance(content_obj, list):
                    line = "".join([p['value'] for p in content_obj if p['type'] == 'text'])
                else:
                    line = str(content_obj)
                line_low = line.lower()

                if line_low.startswith('topic:'):
                    data['topic'] = line.split(':', 1)[1].strip()
                elif line_low.startswith('subtopic:'):
                    data['subtopic'] = line.split(':', 1)[1].strip()
                elif line.strip():
                    data_text_list.append((content_obj, ilvl, is_list))

            if 'topic' not in data and global_metadata.get('topic'):
                data['topic'] = global_metadata['topic']
            if 'subtopic' not in data and global_metadata.get('subtopic'):
                data['subtopic'] = global_metadata['subtopic']

            if data_text_list:
                data['text'] = data_text_list

            apply_metadata_to_slide(slide, data)

            images = section.get('images', [])
            if images:
                for img_blob in images:
                    with Image.open(io.BytesIO(img_blob)) as pil_img:
                        orig_w, orig_h = pil_img.size
                    img_aspect = orig_w / orig_h
                    max_w = int(prs.slide_width * 0.8)
                    max_h = int(prs.slide_height * 0.6)
                    if img_aspect > max_w / max_h:
                        draw_w = max_w
                        draw_h = int(max_w / img_aspect)
                    else:
                        draw_h = max_h
                        draw_w = int(max_h * img_aspect)
                    left = int((prs.slide_width - draw_w) / 2)
                    top = int((prs.slide_height - draw_h) / 2)
                    slide.shapes.add_picture(io.BytesIO(img_blob), left, top, width=draw_w, height=draw_h)

            print(f"Inserted {layout.name} with content")

        # ── science_discussion_page ─────────────────────────────────
        elif layout.name == 'LAYOUT_science_discussion_page':
            # 1. Inject static header ("Discussion Time")
            if discussion_page_data['header_xml'] is not None:
                slide.shapes._spTree.append(copy.deepcopy(discussion_page_data['header_xml']))

            # 2. Inject static subtitle ("Discuss with your partner:")
            if discussion_page_data['subtitle_xml'] is not None:
                slide.shapes._spTree.append(copy.deepcopy(discussion_page_data['subtitle_xml']))

            # 3. Parse discussion questions from DOCX
            discussion_lines = []
            for entry in section['content']:
                line = get_text(entry)
                if line.strip():
                    discussion_lines.append(entry) # Use raw entry for rich text handling

            # 4. Clone the question group template for each question
            q_grp_xml = discussion_page_data['question_grp_xml']
            if q_grp_xml is not None and discussion_lines:
                row_height = int(Inches(1.2)) # spacing between each question
                a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

                # The first question usually starts around Y=3.4 Inches (based on Group 28 top in layout)
                # But we'll just use the template's position and increment.

                for idx, sentence_data in enumerate(discussion_lines):
                    new_grp = copy.deepcopy(q_grp_xml)

                    # Copy image rels if needed
                    if discussion_page_data['layout_part']:
                        copy_image_rels(new_grp, discussion_page_data['layout_part'], slide.part)

                    slide.shapes._spTree.append(new_grp)
                    grp_shape = slide.shapes[-1]

                    # Shift the group down
                    grp_shape.top = grp_shape.top + (idx * row_height)

                    # Update the text in the rounded rectangle
                    for sub in grp_shape.shapes:
                        if getattr(sub, 'has_text_frame', False):
                            sub_text = sub.text.strip().lower()
                            if sub_text == 'text' or len(sub_text) > 1:
                                tf = sub.text_frame
                                p0 = tf.paragraphs[0]
                                for r_elem in list(p0._p.findall(f'{{{a_ns}}}r')):
                                    p0._p.remove(r_elem)
                                
                                # Build content parts
                                parts = get_content_parts(sentence_data)
                                for part in parts:
                                    if part['type'] == 'text':
                                        run = p0.add_run()
                                        run.text = part['value']
                                        run.font.size = Pt(28)
                                        run.font.name = 'Calibri'
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                                    elif part['type'] == 'math':
                                        try:
                                            from lxml.etree import fromstring as lxml_parse
                                            math_elem = lxml_parse(part['value'])
                                            p0._p.append(math_elem)
                                        except Exception as e:
                                            print(f"Error injecting math: {e}")

                print(f"Inserted {layout.name} with {len(discussion_lines)} question items")
            else:
                if not discussion_lines:
                    print(f"WARNING: No content found for {layout.name}")
                else:
                    print(f"WARNING: No question group template found for {layout.name}")

            apply_metadata_to_slide(slide, {})

        # ── finalquiz_page ──────────────────────────────────────────
        elif layout.name == 'LAYOUT_finalquiz_page':
            # Inject static elements for the intro slide
            templates = content_templates.get(layout.name, {})
            for static_elem in templates.get('static_elements', []):
                elem_copy = remove_locks(copy.deepcopy(static_elem))
                copy_image_rels(elem_copy, layout.part, slide.part)
                slide.shapes._spTree.append(elem_copy)
            print("Inserted LAYOUT_finalquiz_page (Intro Slide)")

            # Parse Q/A pairs and generate question + answer slides
            qa_pairs = []
            for entry in section['content']:
                line = get_text(entry)
                if ":" in line:
                    parts = line.split(":", 1)
                    q_text = parts[0].strip()
                    a_text = parts[1].strip()
                    if q_text.lower() == 'question' and a_text.lower() == 'answer':
                        continue
                    qa_pairs.append((q_text, a_text))

            layout_q = None
            layout_a = None
            for lay in prs.slide_layouts:
                if lay.name == 'LAYOUT_finalquiz_page_q':
                    layout_q = lay
                if lay.name == 'LAYOUT_finalquiz_page_a':
                    layout_a = lay

            for q, a in qa_pairs:
                if layout_q:
                    s_q = prs.slides.add_slide(layout_q)
                    inject_logo(s_q)
                    for shape in s_q.shapes:
                        if shape.is_placeholder and shape.placeholder_format.idx == 0:
                            shape.text = q
                if layout_a:
                    s_a = prs.slides.add_slide(layout_a)
                    inject_logo(s_a)
                    for shape in s_a.shapes:
                        if shape.is_placeholder and shape.placeholder_format.idx == 0:
                            shape.text = a
            print(f"Generated {len(qa_pairs)} Question and Answer slide pairs.")

        # ── Default: just apply metadata ───────────────────────────
        else:
            apply_metadata_to_slide(slide, {})
            print(f"Inserted {layout.name} (default handler)")

        # ── SYR overlay ─────────────────────────────────────────────
        if section.get('has_syr'):
            layout_syr = get_layout('syr')
            if layout_syr:
                syr_templates = content_templates.get('LAYOUT_syr', {})
                for static_elem in syr_templates.get('static_elements', []):
                    elem_copy = remove_locks(copy.deepcopy(static_elem))
                    copy_image_rels(elem_copy, layout_syr.part, slide.part)
                    slide.shapes._spTree.append(elem_copy)
                print("Overlaid LAYOUT_syr on current slide.")

        # Inject logo on every slide
        inject_logo(slide)

    report_progress(96, "Saving presentation...")
    prs.save(output_path)
    report_progress(100, "Done!")


if __name__ == "__main__":
    generate_ppt(
        "higher_content.docx",
        "template_higher_class.pptx",
        "Generated_Higher_Presentation.pptx"
    )
    print("DONE -- Higher Class PPT Generated", flush=True)
