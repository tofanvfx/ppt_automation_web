#!/usr/bin/env python3
"""
Worksheet Generator: DOCX → PPTX
==============================================
Converts a structured Word document into a formatted worksheet PowerPoint slide
matching the Aveti Learning template (10.83" x 7.5" landscape, 2-column layout).

Usage:
    python worksheet_generator.py input.docx [output.pptx] [template.pptx]

See sample_worksheet.docx for the required DOCX format.
"""

import sys, re, os
from docx import Document
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ── Colors ──────────────────────────────────────────────────────────────────
C_SECTION = RGBColor(0x00, 0x70, 0xC0)   # Blue  – section headers
C_ANSWER  = RGBColor(0x1F, 0x6B, 0x1F)   # Green – answer text
C_FALSE   = RGBColor(0xFF, 0x00, 0x00)   # Red   – "False" verdict

# ── Layout (EMU; 1 inch = 914400) ────────────────────────────────────────────
CONTENT_TOP = 518483
CONTENT_BTM = 6350000

COL_L_X    = 305530
COL_L_W    = 4540128
COL_R_X    = 5137288
COL_R_W    = 4540128
OPT_INDENT = 210000   # left indent for option lines

FONT = 'Calibri'
FS   = 127000          # 10 pt in EMU

# ── Spacing & Height Calculation ─────────────────────────────────────────────
# We use a unified system: Height = Lines * LINE_H
# Spacing = Height + specific GAP
LINE_H       = 175000  # EMU per line (~10pt with leading)
GAP_INTERNAL = 10000   # Tighter internal gaps
GAP_QUESTION = 20000   # Minimal base gap (mostly rely on DOCX)
GAP_SECTION  = 60000   # Reduced gap after section header
GAP_PARAGRAPH= 100000  # Reduced gap per empty line (~0.5 line)

CHARS_PER_LINE = 76

def text_lines(text: str, chars_per_line: int = CHARS_PER_LINE) -> int:
    """Estimate number of wrapped lines."""
    if not text: return 0
    lines = 0
    for raw_line in text.split('\n'):
        lines += max(1, -(-len(raw_line) // chars_per_line))
    return lines

def get_block_h(b):
    """Calculate just the height of the content box."""
    k = b['kind']
    if k in ('sec_header', 'sec_gap'): return LINE_H if k == 'sec_header' else 0
    if k == 'gap': return 0
    
    if k == 'question':
        if b.get('qtype') == 'AR':
            ass_text = 'Assertion: ' + b.get('assertion', '')
            rea_text = 'Reason: '    + b.get('reason', '')
            n = text_lines(ass_text) + text_lines(rea_text)
            return max(1, n) * LINE_H
        n = text_lines(b['text'])
        return max(1, n) * LINE_H
        
    if k == 'options':
        n = sum(text_lines(opt) for opt in b['options'])
        return max(1, n) * LINE_H
        
    if k == 'answer':
        if b.get('qtype') == 'TF':
            n = text_lines(b['lines'][0])
            return max(1, n) * LINE_H
        n = sum(text_lines(line) for line in b['lines'])
        if len(b['lines']) > 1: n += 1
        return max(1, n) * LINE_H
        
    return 0

def get_block_sp(b):
    """Calculate the Y-advance for the block (Height + Gap)."""
    h = get_block_h(b)
    k = b['kind']
    if k == 'sec_header': return h + GAP_SECTION
    if k == 'sec_gap':    return 40000 
    if k == 'gap':        return GAP_PARAGRAPH
    
    if k == 'question':   return h + GAP_INTERNAL
    if k == 'options':    return h + GAP_INTERNAL
    if k == 'answer':     return h + GAP_QUESTION 
    
    return h

# ─────────────────────────────────────────────────────────────────────────────
# DOCX PARSER
# ─────────────────────────────────────────────────────────────────────────────

def parse_docx(path):
    doc  = Document(path)
    data = dict(chapter='', class_num='', subject='', worksheet='', content=[])
    cur_sec = cur_q = None
    in_ans  = False
    ans_buf = []

    def flush_q():
        nonlocal cur_q, ans_buf, in_ans
        if cur_q:
            cur_q['answer'] = ans_buf[:]
            cur_sec['questions'].append(cur_q)
        cur_q = None; ans_buf = []; in_ans = False

    def flush_sec():
        nonlocal cur_sec
        if cur_sec:
            flush_q()
            data['content'].append(cur_sec)
        cur_sec = None

    for para in doc.paragraphs:
        txt = para.text.strip()
        
        # ── Handle Empty Lines (Gaps) ───────────────────────────────────────
        if not txt:
            # Only record a gap if we are NOT at the very start of a section
            if cur_sec:
                flush_q()
                # Avoid adding multiple gaps or gaps right after section header
                if not cur_sec['questions'] or cur_sec['questions'][-1].get('type') != 'gap':
                    cur_sec['questions'].append(dict(type='gap'))
            else:
                if not data['content'] or data['content'][-1].get('type') != 'gap':
                    data['content'].append(dict(type='gap'))
            continue

        lo = txt.lower()

        # ── Metadata ────────────────────────────────────────────────────────
        matched = False
        for key, field in [('chapter:', 'chapter'), ('class:', 'class_num'),
                            ('subject:', 'subject'), ('worksheet:', 'worksheet')]:
            if lo.startswith(key):
                data[field] = txt.split(':', 1)[1].strip()
                matched = True; break
        if matched:
            continue

        # ── Section header ───────────────────────────────────────────────────
        if re.match(r'^[A-F]\.\s+\S', txt):
            flush_sec()
            cur_sec = dict(title=txt, type=_stype(txt), questions=[])
            continue

        if cur_sec is None:
            continue

        # ── Question ─────────────────────────────────────────────────────────
        if re.match(r'^Q\d+[.\s]', txt):
            flush_q()
            cur_q = dict(text=txt, options=[], answer=[],
                         assertion='', reason='', type=cur_sec['type'])
            m = re.search(r'Assertion:\s*(.+?)(?:\s+Reason:\s*(.+))?$', txt)
            if m:
                cur_q['assertion'] = m.group(1).strip()
                if m.group(2): cur_q['reason'] = m.group(2).strip()
            continue

        if cur_q is None:
            continue

        # ── Assertion / Reason sub-lines ──────────────────────────────────────
        if lo.startswith('assertion:'):
            cur_q['assertion'] = txt.split(':', 1)[1].strip(); in_ans = False; continue
        if lo.startswith('reason:'):
            cur_q['reason']    = txt.split(':', 1)[1].strip(); in_ans = False; continue

        # ── Options ──────────────────────────────────────────────────────────
        if re.match(r'^[a-dA-D][)\.]\s', txt) and not in_ans:
            opt = re.sub(r'^([a-dA-D])[\.]\s', r'\1) ', txt)
            cur_q['options'].append(opt); continue

        # ── Answer ───────────────────────────────────────────────────────────
        if re.match(r'^[Aa]nswer\s*:', txt):
            in_ans = True
            val = re.sub(r'^[Aa]nswer\s*:\s*', '', txt).strip()
            if val: ans_buf.append(val)
            continue
        if in_ans:
            ans_buf.append(re.sub(r'^\d+\.\s*', '', txt))

    flush_sec()
    return data


def _stype(t):
    u = t.upper()
    if 'MULTIPLE' in u or 'MCQ' in u: return 'MCQ'
    if 'FILL'     in u or 'FIB' in u: return 'FIB'
    if 'TRUE'     in u or 'T/F' in u: return 'TF'
    if 'ASSERTION' in u:               return 'AR'
    if 'SHORT'    in u or '(SA)' in u: return 'SA'
    if 'LONG'     in u or '(LA)' in u: return 'LA'
    return 'OTHER'


# ─────────────────────────────────────────────────────────────────────────────
# BLOCK BUILDER  (flat list of renderable units)
# ─────────────────────────────────────────────────────────────────────────────

def build_blocks(content):
    blocks = []
    for item in content:
        if item.get('type') == 'gap':
            blocks.append(dict(kind='gap'))
            continue
        
        # It's a section
        sec = item
        blocks.append(dict(kind='sec_header', text=sec['title']))
        for q in sec['questions']:
            if q.get('type') == 'gap':
                blocks.append(dict(kind='gap'))
                continue
                
            blocks.append(dict(kind='question', text=q['text'],
                               assertion=q.get('assertion',''),
                               reason=q.get('reason',''),
                               qtype=q['type']))

            if q['options']:
                blocks.append(dict(kind='options', options=q['options']))

            if q['answer']:
                blocks.append(dict(kind='answer', lines=q['answer'], qtype=q['type']))
    return blocks


def _bheight(b):  return get_block_h(b)
def _bspacing(b): return get_block_sp(b)


# ─────────────────────────────────────────────────────────────────────────────
# COLUMN FLOW
# ─────────────────────────────────────────────────────────────────────────────

def flow(blocks):
    """Distribute blocks into left and right columns across multiple slides.
    Returns a list of slides, where each slide is (left_placed, right_placed).
    Each placed item is (block, y_position).
    """
    slides = []
    left_placed, right_placed = [], []
    left_y = right_y = CONTENT_TOP
    col = 'left'

    def flush_slide():
        nonlocal left_placed, right_placed, left_y, right_y, col
        slides.append((left_placed, right_placed))
        left_placed, right_placed = [], []
        left_y = right_y = CONTENT_TOP
        col = 'left'

    for i, b in enumerate(blocks):
        sp = _bspacing(b)
        h  = _bheight(b)

        if b['kind'] == 'sec_gap':
            if col == 'left': left_y  += sp
            else:              right_y += sp
            continue

        if col == 'left':
            if left_y + h > CONTENT_BTM:
                col = 'right'
            else:
                # Keep-with-next: if this is a section header, peek at next
                if b['kind'] == 'sec_header':
                    next_h = 0
                    for j in range(i + 1, len(blocks)):
                        nb = blocks[j]
                        if nb['kind'] == 'sec_gap': continue
                        next_h = _bheight(nb)
                        break
                    if next_h and (left_y + sp + next_h > CONTENT_BTM):
                        col = 'right'
                    else:
                        left_placed.append((b, left_y))
                        left_y += sp
                        continue
                else:
                    left_placed.append((b, left_y))
                    left_y += sp
                    continue

        if col == 'right':
            if right_y + h > CONTENT_BTM:
                flush_slide()
                if left_y + h > CONTENT_BTM:
                     left_placed.append((b, left_y))
                     left_y += sp
                else:
                    left_placed.append((b, left_y))
                    left_y += sp
            else:
                right_placed.append((b, right_y))
                right_y += sp

    if left_placed or right_placed:
        slides.append((left_placed, right_placed))

    return slides


# ─────────────────────────────────────────────────────────────────────────────
# RENDERING
# ─────────────────────────────────────────────────────────────────────────────

def _tb(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap     = True
    tf.margin_left   = tf.margin_right  = Emu(0)
    tf.margin_top    = tf.margin_bottom = Emu(0) # Zero margins for absolute gap control
    return tf


def _run(p, text, bold=False, color=None):
    r = p.add_run()
    r.text           = text
    r.font.size      = Emu(FS)
    r.font.bold      = bold
    r.font.name      = FONT
    if color: r.font.color.rgb = color
    return r


def render_block(slide, b, y, cx, cw):
    k = b['kind']
    bh = get_block_h(b)

    if k == 'sec_header':
        tf = _tb(slide, cx, y, cw, bh)
        _run(tf.paragraphs[0], b['text'], bold=True, color=C_SECTION)

    elif k == 'question':
        tf = _tb(slide, cx, y, cw, bh)
        p0 = tf.paragraphs[0]

        if b['qtype'] == 'AR' and (b['assertion'] or b['reason']):
            m = re.match(r'^(Q\d+\.\s*)(.*)', b['text'])
            if m: _run(p0, m.group(1), bold=True)
            _run(p0, 'Assertion: ', bold=True)
            _run(p0, b['assertion'] or (m.group(2) if m else b['text']))
            if b['reason']:
                p1 = tf.add_paragraph()
                _run(p1, 'Reason', bold=True)
                _run(p1, ': ' + b['reason'])
        else:
            # DOTALL used to handle multi-line inputs correctly
            m = re.match(r'^(Q\d+\.\s*)(.*)', b['text'], re.DOTALL)
            if m:
                _run(p0, m.group(1), bold=True)
                _run(p0, m.group(2))
            else:
                _run(p0, b['text'])

    elif k == 'options':
        tf = _tb(slide, cx + OPT_INDENT, y, cw - OPT_INDENT, bh)
        for i, opt in enumerate(b['options']):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _run(p, opt)

    elif k == 'answer':
        lines = b['lines']
        qt    = b['qtype']
        n     = len(lines)

        if qt == 'TF':
            tf = _tb(slide, cx, y, cw, bh)
            p  = tf.paragraphs[0]
            _run(p, 'Answer: ', bold=True, color=C_ANSWER)
            raw = lines[0]
            m   = re.match(r'^(True|False)\s*[–-]\s*(.*)', raw, re.DOTALL)
            if m:
                v, e = m.group(1), m.group(2)
                _run(p, v, bold=True, color=C_FALSE if v == 'False' else C_ANSWER)
                _run(p, ' \u2013 ' + e, color=C_ANSWER)
            else:
                _run(p, raw, color=C_ANSWER)

        elif n == 1:
            lbl = '✔ Correct Answer: ' if qt in ('MCQ', 'AR') else '✔ Answer: '
            tf  = _tb(slide, cx, y, cw, bh)
            p   = tf.paragraphs[0]
            _run(p, lbl, bold=True, color=C_ANSWER)
            _run(p, lines[0], color=C_ANSWER)

        else:
            tf = _tb(slide, cx, y, cw, bh)
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                if i == 0: _run(p, '✔ Answer:', bold=True, color=C_ANSWER)
                else:       _run(p, line,         color=C_ANSWER)


def render_column(slide, placed, cx, cw):
    for b, y in placed:
        render_block(slide, b, y, cx, cw)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE MANAGEMENT
# ─────────────────────────────────────────────────────────────────────────────

def copy_header_footer(source_slide, target_slide):
    """Deep copy the metadata text boxes (shapes 0-3) from source to target."""
    for i in range(4): # Only Chapter, Class, Subject, Worksheet
        if i >= len(source_slide.shapes): break
        shape = source_slide.shapes[i]
        # Copy the shape
        import copy
        el = shape.element
        new_el = copy.deepcopy(el)
        target_slide.shapes._spTree.append(new_el)

def update_header(slide, data):
    shapes = list(slide.shapes)
    # Mapping based on inspection: 0: Worksheet, 1: Class, 2: Subject, 3: Chapter
    upd = {
        0: data['worksheet'],
        1: f"Class {data['class_num']}" if data['class_num'] else '',
        2: data['subject'],
        3: data['chapter'],
    }
    for idx, text in upd.items():
        if not text: continue
        if idx >= len(shapes): continue
        sp = shapes[idx]
        if not hasattr(sp, 'text_frame'): continue
        tf = sp.text_frame
        p  = tf.paragraphs[0]
        
        # Preserve existing formatting if possible
        bold_v = size_v = col_v = None
        if p.runs:
            r0 = p.runs[0]
            bold_v = r0.font.bold
            size_v = r0.font.size
            try: col_v = r0.font.color.rgb
            except Exception: pass
        
        # Clear and add new run
        for el in list(p._p.findall(qn('a:r'))):
            p._p.remove(el)
        r = p.add_run()
        r.text = text
        r.font.bold = bold_v; r.font.size = size_v; r.font.name = FONT
        if col_v: r.font.color.rgb = col_v


def remove_content_shapes(slide):
    """Remove shapes from Slide 1 starting from index 4 (content begins)."""
    tree   = slide.shapes._spTree
    shapes = list(slide.shapes)
    # Indices 0-3 are metadata (preserved). Index 4+ are content.
    for s in shapes[4:]:
        tree.remove(s._element)


# ─────────────────────────────────────────────────────────────────────────────
# ENTRY POINT
# ─────────────────────────────────────────────────────────────────────────────

def generate(docx_path, template_pptx, output_pptx, progress_callback=None):
    def report_progress(percentage, status):
        if progress_callback:
            progress_callback(percentage, status)

    report_progress(5, "Parsing DOCX...")
    print(f"\U0001f4c4 Parsing DOCX: {docx_path}")
    data = parse_docx(docx_path)
    print(f"   Chapter  : {data['chapter']}")
    print(f"   Class    : {data['class_num']}")
    print(f"   Subject  : {data['subject']}")
    print(f"   Worksheet: {data['worksheet']}")
    num_secs = sum(1 for item in data['content'] if item.get('type') != 'gap')
    print(f"   Sections : {num_secs}")

    report_progress(20, "Loading template...")
    print(f"\n\U0001f4d0 Loading template: {template_pptx}")
    prs   = Presentation(template_pptx)
    template_slide = prs.slides[0]

    report_progress(30, "Building & flowing blocks...")
    print("\U0001f4ca Building & flowing blocks…")
    blocks = build_blocks(data['content'])
    slide_data = flow(blocks)
    print(f"   {len(blocks)} blocks → {len(slide_data)} slides")

    # Use the template slide's layout
    target_layout = template_slide.slide_layout

    total_slides = len(slide_data)
    for i, (left_placed, right_placed) in enumerate(slide_data):
        slide_prog = 30 + int(((i) / total_slides) * 65)
        report_progress(slide_prog, f"Rendering slide {i+1} of {total_slides}...")
        print(f"   Rendering slide {i+1}...")
        if i == 0:
            slide = template_slide
            remove_content_shapes(slide)
        else:
            slide = prs.slides.add_slide(target_layout)
            copy_header_footer(template_slide, slide)

        print(f"     Updating header on slide {i+1}…")
        update_header(slide, data)

        print(f"     Rendering columns on slide {i+1}…")
        render_column(slide, left_placed,  COL_L_X, COL_L_W)
        render_column(slide, right_placed, COL_R_X, COL_R_W)

    try:
        report_progress(95, "Saving PPTX...")
        prs.save(output_pptx)
        report_progress(100, "Completed")
        print(f"\n\u2705 Saved \u2192 {output_pptx}")
    except PermissionError:
        print(f"\n\u274c Error: Could not save to {output_pptx}.")
        print("   Please ensure the PowerPoint file is CLOSED and try again.")
        sys.exit(1)


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python worksheet_generator.py input.docx [output.pptx] [template.pptx]")
        sys.exit(0)

    docx_in  = sys.argv[1]
    pptx_out = sys.argv[2] if len(sys.argv) > 2 else docx_in.replace('.docx', '_worksheet.pptx')
    template = sys.argv[3] if len(sys.argv) > 3 else \
               os.path.join(os.path.dirname(os.path.abspath(__file__)), 'worksheet_test.pptx')

    for f, lbl in [(docx_in, 'DOCX'), (template, 'Template PPTX')]:
        if not os.path.exists(f):
            print(f"Error: {lbl} not found: {f}"); sys.exit(1)

    generate(docx_in, template, pptx_out)
