import re
import os
import copy
import io
from docx import Document
from docx.document import Document as _Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml import parse_xml
from PIL import Image


def generate_ppt(docx_path, template_path, output_path):
    prs = Presentation(template_path)
    doc = Document(docx_path)

    def clean(text):
        return text.strip().lower()

    def get_text(entry):
        """Extract plain text from a content entry (either a string or a (text_parts, ilvl) tuple)."""
        data = entry[0] if isinstance(entry, tuple) else entry
        if isinstance(data, list):
            # Concatenate only text parts for plain text version
            return "".join(p['value'] if p['type'] == 'text' else '' for p in data)
        return str(data)

    def copy_image_rels(element, source_part, target_part):
        """Copy image relationships referenced by blip elements from source to target part."""
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        r_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
        for blip in element.findall('.//a:blip', namespaces=ns):
            old_rid = blip.get(f'{r_ns}embed')
            if old_rid and old_rid in source_part.rels:
                rel = source_part.rels[old_rid]
                # Add the image to the target part and get a new rId
                new_rid = target_part.relate_to(rel.target_part, rel.reltype)
                blip.set(f'{r_ns}embed', new_rid)

    def remove_locks(element):
        """Remove a:spLocks and a:grpSpLocks elements from a shape XML to make it fully editable."""
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        for lock in element.findall('.//a:spLocks', namespaces=ns):
            lock.getparent().remove(lock)
        for lock in element.findall('.//a:grpSpLocks', namespaces=ns):
            lock.getparent().remove(lock)
        return element

    def get_layout(name):
        for layout in prs.slide_layouts:
            if layout.name == f"LAYOUT_{name}":
                return layout
            
        # Try case-insensitive, underscore-insensitive match
        target = f"layout_{name}".replace("_", "").lower()
        for layout in prs.slide_layouts:
            if layout.name.replace("_", "").replace(" ", "").lower() == target:
                return layout
            
        if name.replace("_", "").lower() == 'mathtitlepage':
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_math_page_title':
                    return layout

        if name.replace("_", "").lower() == 'ssttitlepage':
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_page_title':
                    return layout
                
        if name.replace("_", "").lower() in ('learningobjective', 'sstlopage', '1layoutsstlopage'):
            for layout in prs.slide_layouts:
                if layout.name == '1_LAYOUT_sst_lo_page' or layout.name == 'LAYOUT_sst_lo_page':
                    return layout

        if name.replace("_", "").lower() == 'mathlopage':
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_math_lo_page':
                    return layout

        if name.replace("_", "").lower() in ('finalquizpage', 'finalquiz', 'quiz'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_final_quiz_page':
                    return layout

        if name.replace("_", "").lower() in ('quiztimepage', 'sstquiztimepage', 'sst_quiztime_page'):
            # Return _01 as default; the main loop will override based on option lengths
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_quiztime_page_01':
                    return layout

        if name.replace("_", "").lower() in ('sstcontentpage01', 'sstcontentpage1'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_content_page_01':
                    return layout

        if name.replace("_", "").lower() in ('sstsummarypage', 'sstsummary', 'sstsummerypage', 'sstsummery'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_summary_page':
                    return layout

        if name.replace("_", "").lower() in ('mathsummarypage', 'mathsummary', 'mathsummerypage', 'mathsummery'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_math_summary_page':
                    return layout

        if name.replace("_", "").lower() in ('notedownpage', 'sstnotedownpage', 'sst_notedown_page'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_notedown_page':
                    return layout

        if name.replace("_", "").lower() in ('previouspage', 'sstpreviouspage', 'sst_previous_page'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_previous_page':
                    return layout

        if name.replace("_", "").lower() in ('homeworkpage', 'ssthomeworkpage', 'ssthomework', 'homework'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_homework_page':
                    return layout

        if name.replace("_", "").lower() in ('discussionpage', 'sstdiscussionpage', 'sstdiscussion', 'sst_discussion_page'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_discussion_page':
                    return layout

        if name.replace("_", "").lower() in ('mathdefaultpage', 'math_default_page'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_math_default_page':
                    return layout

        if name.replace("_", "").lower() in ('sstdeafultpage', 'sst_deafult_page'):
            for layout in prs.slide_layouts:
                if layout.name == 'LAYOUT_sst_deafult_page':
                    return layout

        return None

    # Nested dictionary to store templates per layout
    sst_content_templates = {
        'LAYOUT_sst_content_page_01': {'topic': None, 'subtopic': None, 'text': None},
        'LAYOUT_sst_content_page_02': {'topic': None, 'subtopic': None, 'text': None},
        'LAYOUT_sst_notedown_page': {'text': None},
        'LAYOUT_sst_quiztime_page_01': {'title': None, 'question': None, 'options': [], 'picture': None},
        'LAYOUT_sst_quiztime_page_02': {'title': None, 'question': None, 'options': [], 'picture': None},
        'LAYOUT_sst_discussion_page': {'question1': None, 'static_elements': []},
        'LAYOUT_sst_homework_page': {'static_elements': []},
        'LAYOUT_syr': {'static_elements': []},
        'LAYOUT_ask_question': {'static_elements': []},
        'LAYOUT_sst_activity_page_01': {'static_elements': [], 'text': None},
        'LAYOUT_sst_activity_page_02': {'static_elements': [], 'text': None},
        'LAYOUT_sst_deafult_page': {'topic': None, 'subtopic': None, 'text': None},
        'LAYOUT_math_default_page': {'topic': None, 'subtopic': None, 'text': None}
    }

    # Extract lo_page group shape XML before processing slides
    lo_group_xmls = {
        '1_LAYOUT_sst_lo_page': None,
        'LAYOUT_sst_lo_page': None,
        'LAYOUT_math_lo_page': None,
        'LAYOUT_sst_summary_page': None,
        'LAYOUT_math_summary_page': None,
        'LAYOUT_sst_previous_page': None
    }

    # Store static title templates for lo/summary layouts where add_slide fails
    lo_title_xmls = {k: None for k in lo_group_xmls.keys()}
    lo_title_xmls['1_LAYOUT_sst_lo_page'] = None # Ensure renamed one is included
    lo_title_xmls['LAYOUT_sst_previous_page'] = None

    for layout in prs.slide_layouts:
        if layout.name in lo_group_xmls:
            for shape in layout.shapes:
                if shape.shape_type == 6:  # msoGroup
                    for subshape in shape.shapes:
                        if getattr(subshape, 'has_text_frame', False) and 'Text goes here' in subshape.text:
                            lo_group_xmls[layout.name] = copy.deepcopy(shape.element)
                            shape.element.getparent().remove(shape.element)
                            break
                    if lo_group_xmls[layout.name] is not None:
                        break
    
        if layout.name in lo_title_xmls:
            for shape in layout.shapes:
                if shape.shape_type != 6 and getattr(shape, 'has_text_frame', False) and shape.text.strip():
                    lo_title_xmls[layout.name] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                    # print(f"DEBUG: Captured title for {layout.name}: '{shape.text[:20]}...'")
                    break
    
        if layout.name in sst_content_templates:
            templates = sst_content_templates[layout.name]
            for shape in list(layout.shapes): # Use list() to avoid issues when removing
                has_topic = False
                has_subtopic = False
                has_quiz_title = False
            
                # 0. Capture standalone PICTURE shapes for quiztime layouts
                if layout.name.startswith('LAYOUT_sst_quiztime_page') and shape.shape_type == 13:  # PICTURE
                    templates['picture'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                    continue

                # 0b. Capture ALL non-question1 shapes for discussion layout (groups, pictures, etc.)
                if layout.name == 'LAYOUT_sst_discussion_page':
                    if not (shape.has_text_frame and 'question1' in shape.text.lower()):
                        templates['static_elements'].append(copy.deepcopy(shape.element))
                        shape.element.getparent().remove(shape.element)
                        continue

                # 0c. Capture ALL shapes for homework layout
                if layout.name == 'LAYOUT_sst_homework_page':
                    templates['static_elements'].append(copy.deepcopy(shape.element))
                    shape.element.getparent().remove(shape.element)
                    continue

                # 0d. Capture ALL shapes for syr layout
                if layout.name == 'LAYOUT_syr':
                    templates['static_elements'].append(copy.deepcopy(shape.element))
                    shape.element.getparent().remove(shape.element)
                    continue

                # 0e. Capture ALL shapes for ask_question layout
                if layout.name == 'LAYOUT_ask_question':
                    templates['static_elements'].append(copy.deepcopy(shape.element))
                    shape.element.getparent().remove(shape.element)
                    continue

                # 0f. For activity pages, separate static elements from the primary text box
                if layout.name in ('LAYOUT_sst_activity_page_01', 'LAYOUT_sst_activity_page_02'):
                    # Exclude Picture Placeholders from extraction so they stay in layout
                    if shape.is_placeholder and shape.placeholder_format.type == 18:
                        continue

                    if not (shape.has_text_frame and 'text goes here' in shape.text.lower()):
                        templates['static_elements'].append(copy.deepcopy(shape.element))
                        shape.element.getparent().remove(shape.element)
                        continue

                # 1. Check groups for topic/subtopic or "QUIZ TIME"
                if shape.shape_type == 6:  # GROUP
                    for sub in shape.shapes:
                        if getattr(sub, 'has_text_frame', False):
                            txt = sub.text.strip().lower()
                            if txt == 'topic': has_topic = True
                            elif txt == 'subtopic': has_subtopic = True
                            elif layout.name.startswith('LAYOUT_sst_quiztime_page') and 'quiz time' in txt:
                                has_quiz_title = True
                
                    if has_quiz_title:
                        templates['title'] = copy.deepcopy(shape.element)
                        shape.element.getparent().remove(shape.element)
                        continue # Move to next shape

                # 2. Check individual shapes for topic/subtopic
                elif getattr(shape, 'has_text_frame', False):
                    txt = shape.text.strip().lower()
                    if txt == 'topic': has_topic = True
                    elif txt == 'subtopic': has_subtopic = True
            
                if has_topic:
                    templates['topic'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                elif has_subtopic:
                    templates['subtopic'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
            
                # 3. Check for body/placeholder text (question, options, or general text)
                elif shape.has_text_frame and ('Text goes here' in shape.text or (shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY) or (layout.name == 'LAYOUT_sst_discussion_page' and 'question1' in shape.text.lower())):
                    txt = shape.text.lower()
                    if layout.name.startswith('LAYOUT_sst_quiztime_page'):
                        if 'quiz time' in txt:
                            templates['title'] = copy.deepcopy(shape.element)
                        elif 'question' in txt:
                            templates['question'] = copy.deepcopy(shape.element)
                        elif 'options' in txt:
                            templates['options'].append(copy.deepcopy(shape.element))
                        elif txt.strip(): 
                            if not templates['title']: templates['title'] = copy.deepcopy(shape.element)
                    elif layout.name == 'LAYOUT_sst_discussion_page':
                        if 'question1' in txt:
                            templates['question1'] = copy.deepcopy(shape.element)
                    else:
                        templates['text'] = copy.deepcopy(shape.element)
                
                    if shape.element.getparent() is not None:
                        shape.element.getparent().remove(shape.element)

    def iter_block_items(parent):
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            raise ValueError("iter_block_items: parent must be Document or _Cell")

        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)


    # Parse DOCX into sections
    sections = []
    current_section = None

    w_ns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'

    for item in iter_block_items(doc):
        lines_to_process = []
        images_to_add = []
    
        if isinstance(item, Paragraph):
            # Detect indent level from numPr/ilvl
            ilvl = 0
            numPr = item._element.find(f'.//{{{w_ns}}}numPr')
            if numPr is not None:
                ilvl_elem = numPr.find(f'{{{w_ns}}}ilvl')
                if ilvl_elem is not None:
                    ilvl = int(ilvl_elem.get(f'{{{w_ns}}}val', '0'))
        
            # Check for math equations (OMML)
            m_ns = {'m': 'http://schemas.openxmlformats.org/officeDocument/2006/math'}
            math_elements = item._element.findall('.//m:oMath', namespaces=m_ns)
        
            if math_elements:
                # If paragraph has math, decompose into parts
                parts = []
                # We iterate through all children of the paragraph element
                for child in item._element.iterchildren():
                    if child.tag.endswith('}oMath'):
                        from lxml import etree
                        math_xml = etree.tostring(child, encoding='unicode')
                        parts.append({'type': 'math', 'value': math_xml})
                    elif child.tag.endswith('}oMathPara'):
                        # oMathPara contains one or more oMath elements
                        for subchild in child.iterchildren():
                            if subchild.tag.endswith('}oMath'):
                                from lxml import etree
                                math_xml = etree.tostring(subchild, encoding='unicode')
                                parts.append({'type': 'math', 'value': math_xml})
                    elif child.tag.endswith('}r'):
                        t_elem = child.find(f'.//{{{w_ns}}}t')
                        if t_elem is not None and t_elem.text:
                            parts.append({'type': 'text', 'value': t_elem.text})
            
                # If parts is empty (sometimes iterchildren misses things or structure is complex),
                # fall back to plain text
                if not parts:
                    lines_to_process.append((item.text.strip(), ilvl))
                else:
                    lines_to_process.append((parts, ilvl))
            else:
                lines_to_process.append((item.text.strip(), ilvl))
            
            # Extract images from paragraph
            for run in item.runs:
                for drawing in run._element.findall('.//w:drawing', namespaces=run._element.nsmap):
                    for blip in drawing.findall('.//a:blip', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                        embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if embed:
                            blob = doc.part.related_parts[embed].blob
                            images_to_add.append(blob)
        elif isinstance(item, Table):
            for row in item.rows:
                # Simple conversion: join first two cells if they exist with ":"
                cells = [c.text.strip() for c in row.cells]
                if len(cells) >= 2:
                    key = cells[0].rstrip(':').strip()
                    val = cells[1].lstrip(':').strip()
                    lines_to_process.append((f"{key}: {val}", 0))
                elif len(cells) == 1:
                    lines_to_process.append((cells[0], 0))

        for entry in lines_to_process:
            content_data = entry[0] if isinstance(entry, tuple) else entry
            ilvl = entry[1] if isinstance(entry, tuple) else 0
        
            if not content_data and not images_to_add:
                continue
            
            # Get plain text for regex matching
            if isinstance(content_data, list):
                search_text = "".join(p['value'] if p['type'] == 'text' else '' for p in content_data)
            else:
                search_text = str(content_data)

            # Extract the first bracketed text as the layout name (e.g., [LAYOUT_sst_content_page_01])
            match = re.search(r'^\[\s*([^\]]+?)\s*\]', search_text)
            if match:
                # Check if there's a SYR marker on the SAME line as the layout marker
                lower_text = search_text.lower()
                has_syr = '[add_syr]' in lower_text or '[syr]' in lower_text
            
                # Check for ask_question marker on the SAME line
                aq_match = re.search(r'\[(?:add_question|ask_question)\s*\((.*?)\)\]', search_text, flags=re.IGNORECASE)
                ask_question_text = aq_match.group(1).strip() if aq_match else None
            
                name = match.group(1).strip()
            
                current_section = {
                    'name': name,
                    'content': [],
                    'images': [],
                    'has_syr': has_syr,
                    'ask_question_text': ask_question_text
                }
                sections.append(current_section)
            elif current_section is not None:
                if content_data:
                    # Use search_text for marker detection/removal
                    lower_text = search_text.lower()
                    if '[add_syr]' in lower_text or '[syr]' in lower_text:
                        current_section['has_syr'] = True
                        # If it's a string, remove it. If it's a list, it's complex, 
                        # but usually markers are at the start/end of a paragraph.
                        if isinstance(content_data, str):
                            content_data = re.sub(r'\[add_syr\]|\[syr\]', '', content_data, flags=re.IGNORECASE).strip()
                
                    # Check for ask_question marker
                    aq_match = re.search(r'\[(?:add_question|ask_question)\s*\((.*?)\)\]', search_text, flags=re.IGNORECASE)
                    if aq_match:
                        current_section['ask_question_text'] = aq_match.group(1).strip()
                        if isinstance(content_data, str):
                            content_data = re.sub(r'\[(?:add_question|ask_question)\s*\(.*?\)\]', '', content_data, flags=re.IGNORECASE).strip()

                    if content_data:
                        current_section['content'].append((content_data, ilvl))
                if images_to_add:
                    current_section['images'].extend(images_to_add)
                    images_to_add = []

    def replace_text_preserve_format(shape, new_text, center=False, font_color=None):
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        if not tf.paragraphs or not tf.paragraphs[0].runs:
            if isinstance(new_text, list):
                shape.text = "\n".join(new_text)
            else:
                shape.text = str(new_text)
            if center:
                for p in tf.paragraphs: p.alignment = PP_ALIGN.CENTER
            return
        
        p0 = tf.paragraphs[0]
        p0_xml = copy.deepcopy(p0._p)
    
        font = p0.runs[0].font
        name = font.name
        size = font.size
        bold = font.bold
        italic = font.italic
        underline = font.underline

        try:
            color_type = getattr(font.color, 'type', None)
            color_rgb = getattr(font.color, 'rgb', None) if color_type == 1 else None
            color_theme = getattr(font.color, 'theme_color', None) if color_type == 2 else None
        except:
            color_type, color_rgb, color_theme = None, None, None

        tf.clear()
    
        texts = new_text if isinstance(new_text, list) else [str(new_text)]
    
        for i, para_data in enumerate(texts):
            if i == 0:
                p = tf.paragraphs[0]
                new_p_xml = copy.deepcopy(p0_xml)
                p._p.getparent().replace(p._p, new_p_xml)
                p = tf.paragraphs[0]
            else:
                new_p_xml = copy.deepcopy(p0_xml)
                tf._txBody.append(new_p_xml)
                p = tf.paragraphs[i]
                p.space_before = Pt(12)
            
            # Aggressively clear existing runs and fields in the paragraph XML
            for r_elem in p._p.findall('.//a:r', namespaces=p._p.nsmap):
                p._p.remove(r_elem)
            for fld_elem in p._p.findall('.//a:fld', namespaces=p._p.nsmap):
                p._p.remove(fld_elem)
            
            if center:
                p.alignment = PP_ALIGN.CENTER
            
            # para_data can be a string OR a list of parts [{'type': 'text', 'value': '...'}, {'type': 'math', 'value': '...'}]
            parts = para_data if isinstance(para_data, list) else [{'type': 'text', 'value': str(para_data)}]
        
            for part in parts:
                if part['type'] == 'text':
                    new_run = p.add_run()
                    new_run.text = part['value']
                
                    if name is not None: new_run.font.name = name
                    if size is not None: new_run.font.size = size
                    if bold is not None: new_run.font.bold = bold
                    if italic is not None: new_run.font.italic = italic
                    if underline is not None: new_run.font.underline = underline
                
                    if font_color is not None:
                        new_run.font.color.rgb = font_color
                    else:
                        if color_rgb is not None:
                            new_run.font.color.rgb = color_rgb
                        elif color_theme is not None:
                            new_run.font.color.theme_color = color_theme
                elif part['type'] == 'math':
                    # Inject OMML XML
                    try:
                        # The OMML might need to be cleaned up or namespaces handled
                        math_elem = parse_xml(part['value'])
                        p._p.append(math_elem)
                    except Exception as e:
                        print(f"Error injecting math XML: {e}")

    for section in sections:
        sname = section['name'].strip().lower()
        if sname == 'sst_content_page':
            if len(section.get('images', [])) > 1:
                layout = get_layout('sst_content_page_02')
            else:
                layout = get_layout('sst_content_page_01')
        elif sname in ('sst_quiztime_page', 'quiztime_page', 'sstquiztimepage'):
            # Parse quiz content first to determine layout
            quiz_data = {'question': '', 'options': []}
            for entry in section['content']:
                line = get_text(entry)
                line_low = line.lower()
                if 'question:' in line_low:
                    quiz_data['question'] = line.split(':', 1)[1].strip()
                elif 'options:' in line_low:
                    opts = line.split(':', 1)[1].strip()
                    quiz_data['options'] = [o.strip() for o in opts.split(',')]
                elif line.strip():
                    if not quiz_data['question']: quiz_data['question'] = line.strip()
                    else: quiz_data['options'].append(line.strip())
        
            # Pick layout based on option character length
            # If ANY option > 25 chars → _01 (long), else _02 (short)
            use_long = any(len(opt) > 25 for opt in quiz_data['options'])
            if use_long:
                layout = get_layout('sst_quiztime_page_01')
                print(f"Quiz options have long text (>25 chars) -> using LAYOUT_sst_quiztime_page_01")
            else:
                layout = get_layout('sst_quiztime_page_02')
                print(f"Quiz options are short (<=25 chars) -> using LAYOUT_sst_quiztime_page_02")
        elif sname in ('sst_activity_page', 'activity_page', 'activity'):
            if len(section.get('images', [])) > 1:
                layout = get_layout('sst_activity_page_01')
            elif len(section['content']) > 1:
                layout = get_layout('sst_activity_page_02')
            else:
                layout = get_layout('sst_activity_page_01')
        else:
            layout = get_layout(section['name'])

        if not layout:
            print(f"Skipping section [{section['name']}], layout not found.")
            continue
        
        slide = prs.slides.add_slide(layout)
    
        if section['name'].replace("_", "").lower() in ('mathpagetitle', 'mathtitlepage', 'sstpagetitle', 'ssttitlepage'):
            data = {}
            for entry in section['content']:
                line = get_text(entry)
                if ":" in line:
                    parts = line.split(":", 1)
                    data[parts[0].strip().upper()] = parts[1].strip()
        
            idx_mapping = {}
            for shape in layout.shapes:
                if shape.is_placeholder and shape.has_text_frame:
                    idx_mapping[clean(shape.text)] = shape.placeholder_format.idx
                
            mapping = {
                "CLASS": "class",
                "SUBJECT": "subject",
                "CHAPTER_NUMBER": "chapter number",
                "CHAPTER_NAME": "chapter name",
                "LESSON": "lesson",
                "TOPIC": "topic"
            }

            for key, template_word in mapping.items():
                if key in data:
                    cleaned_word = clean(template_word)
                    if cleaned_word in idx_mapping:
                        idx = idx_mapping[cleaned_word]
                        found_on_slide = False
                        for shape in slide.shapes:
                            if shape.is_placeholder and shape.placeholder_format.idx == idx:
                                shape.text = data[key]
                                print(f"Updated PageTitle: {template_word} -> {data[key]}")
                                found_on_slide = True
                                break
                        if not found_on_slide:
                            print(f"[{template_word}] idx {idx} NOT FOUND on slide shapes!")
                    else:
                        print(f"[{template_word}] NOT FOUND in idx_mapping! keys={idx_mapping.keys()}")
                else:
                    print(f"[{key}] NOT FOUND in docx data!")
        elif layout.name in lo_group_xmls:
            # Inject static title if available
            if layout.name in lo_title_xmls and lo_title_xmls[layout.name] is not None:
                slide.shapes._spTree.append(copy.deepcopy(lo_title_xmls[layout.name]))

            group_xml = lo_group_xmls.get(layout.name)
            
            if group_xml is not None:
                # Group content: merge sub-bullets (ilvl > 0) with their parent line
                # Each item: (main_text, [(sub_text, ilvl), ...])
                grouped_items = []
                for entry in section['content']:
                    text = entry[0] if isinstance(entry, tuple) else entry
                    ilvl = entry[1] if isinstance(entry, tuple) else 0
                    clean_line = text.strip()
                    if clean_line.startswith(('•', '-', '*')):
                        clean_line = clean_line.lstrip('•-*').strip()
                
                    if ilvl == 0:
                        grouped_items.append((clean_line, []))
                    elif ilvl > 0 and grouped_items:
                        grouped_items[-1][1].append((clean_line, ilvl))
                    else:
                        grouped_items.append((clean_line, []))
            
                current_top_offset = 0
                for main_text, sub_entries in grouped_items:
                    new_element = copy.deepcopy(group_xml)
                    slide.shapes._spTree.append(new_element)
                    new_shape = slide.shapes[-1]
                    new_shape.top = new_shape.top + current_top_offset
                
                    for subshape in new_shape.shapes:
                        if getattr(subshape, 'has_text_frame', False) and 'Text goes here' in subshape.text:
                            tf = subshape.text_frame
                            # Set main text in first paragraph
                            tf.paragraphs[0].text = main_text
                            for run in tf.paragraphs[0].runs:
                                run.font.size = Pt(36)
                                run.font.color.rgb = RGBColor(255, 255, 255)
                        
                            # Add sub-bullet paragraphs with nested indentation
                            for bullet_text, bullet_ilvl in sub_entries:
                                p = tf.add_paragraph()
                                # Indent and bullet marker based on nesting level
                                indent = '    ' * bullet_ilvl
                                markers = {1: '•', 2: '◦', 3: '▪'}
                                marker = markers.get(bullet_ilvl, '▪')
                                p.text = f'{indent}{marker} {bullet_text}'
                                p.space_before = Pt(4)
                                # Smaller font for deeper levels
                                font_size = max(20, 36 - (bullet_ilvl * 4))
                                for run in p.runs:
                                    run.font.size = Pt(font_size)
                                    run.font.color.rgb = RGBColor(255, 255, 255)
                
                    # Calculate height based on content
                    chars_per_line = 65
                    line_count = max(1, len(main_text) // chars_per_line + (1 if len(main_text) % chars_per_line > 0 else 0))
                    sub_height = len(sub_entries) * 0.45
                    current_top_offset += int(Inches(0.9 + 0.55 * (line_count - 1) + sub_height))
                
                total_subs = sum(len(s) for _, s in grouped_items)
                print(f"Updated {layout.name} with {len(grouped_items)} items ({total_subs} sub-bullets)")
        elif layout.name in ('LAYOUT_sst_quiztime_page_01', 'LAYOUT_sst_quiztime_page_02'):
            templates = sst_content_templates.get(layout.name, {})
        
            # Inject static picture
            if templates.get('picture') is not None:
                pic_elem = copy.deepcopy(templates['picture'])
                copy_image_rels(pic_elem, layout.part, slide.part)
                slide.shapes._spTree.append(pic_elem)

            # Inject static title
            if templates.get('title') is not None:
                title_elem = copy.deepcopy(templates['title'])
                copy_image_rels(title_elem, layout.part, slide.part)
                slide.shapes._spTree.append(title_elem)

            # Reuse quiz_data if already parsed during layout selection, otherwise parse now
            if 'quiz_data' not in dir():
                quiz_data = {'question': '', 'options': []}
                for entry in section['content']:
                    content_obj = entry[0] if isinstance(entry, tuple) else entry
                    if isinstance(content_obj, list):
                        # For quiz, we usually want text but can support math in question/options
                        # For now, append the rich object to whichever field is current
                        if not quiz_data['question']: quiz_data['question'] = content_obj
                        else: quiz_data['options'].append(content_obj)
                        continue

                    line = str(content_obj)
                    line_low = line.lower()
                    if 'question:' in line_low:
                        quiz_data['question'] = line.split(':', 1)[1].strip()
                    elif 'options:' in line_low:
                        opts = line.split(':', 1)[1].strip()
                        quiz_data['options'] = [o.strip() for o in opts.split(',')]
                    elif line.strip():
                        if not quiz_data['question']: quiz_data['question'] = line.strip()
                        else: quiz_data['options'].append(line.strip())

            if templates.get('question') is not None:
                question_elem = copy.deepcopy(templates['question'])
            
                # If there are no options, vertically and horizontally center the question text
                if not quiz_data['options']:
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    off = question_elem.find('.//a:off', namespaces=ns)
                    ext = question_elem.find('.//a:ext', namespaces=ns)
                    if off is not None and ext is not None:
                        slide_height = prs.slide_height
                        slide_width = prs.slide_width
                        shape_height = int(ext.get('cy', '0'))
                        shape_width = int(ext.get('cx', '0'))
                    
                        new_y = int((slide_height - shape_height) / 2)
                        new_x = int((slide_width - shape_width) / 2)
                    
                        off.set('y', str(new_y))
                        off.set('x', str(new_x))
            
                slide.shapes._spTree.append(question_elem)
                replace_text_preserve_format(slide.shapes[-1], quiz_data['question'])
        
            # Inject options only if they exist
            if quiz_data['options']:
                option_templates = templates.get('options', [])
                if len(option_templates) == 1:
                    # _01 style: single placeholder, all options as multi-line text
                    slide.shapes._spTree.append(copy.deepcopy(option_templates[0]))
                    replace_text_preserve_format(slide.shapes[-1], quiz_data['options'])
                elif len(option_templates) > 1:
                    # _02 style: separate placeholder per option
                    for idx, opt_template in enumerate(option_templates):
                        slide.shapes._spTree.append(copy.deepcopy(opt_template))
                        if idx < len(quiz_data['options']):
                            replace_text_preserve_format(slide.shapes[-1], quiz_data['options'][idx])
                        else:
                            replace_text_preserve_format(slide.shapes[-1], '')
        
            print(f"Updated {layout.name} with quiz question and {len(quiz_data['options'])} options.")
            quiz_data = None  # Reset for next quiz section
        elif layout.name in ('LAYOUT_sst_content_page_01', 'LAYOUT_sst_content_page_02', 'LAYOUT_sst_deafult_page', 'LAYOUT_math_default_page'):
            data = {}
            data_text_list = []
            for entry in section['content']:
                content_obj = entry[0] if isinstance(entry, tuple) else entry
            
                # If it's a math/multi-part paragraph, it's always body text
                if isinstance(content_obj, list):
                    data_text_list.append(content_obj)
                    continue
                
                line = str(content_obj)
                line_low = line.lower()
            
                if line_low.startswith('topic:'):
                    data['topic'] = line.split(':', 1)[1].strip()
                elif line_low.startswith('subtopic:'):
                    data['subtopic'] = line.split(':', 1)[1].strip()
                elif ":" in line:
                    # Check if it's a known metadata field or just text with a colon
                    parts = line.split(":", 1)
                    key = parts[0].strip().lower()
                    val = parts[1].strip()
                    if key == 'text':
                        data_text_list.append(val)
                    elif key in ('topic', 'subtopic'):
                        data[key] = val
                    else:
                        # Treat unknown "key: val" as plain text if it looks like a sentence
                        data_text_list.append(line)
                elif line.strip():
                    data_text_list.append(line.strip())
                
            if data_text_list:
                data['text'] = data_text_list
        
            # Inject extracted XML templates onto this slide from the correct layout version
            shape_elements = []
            templates = sst_content_templates.get(layout.name, {})
        
            if templates.get('topic') is not None:
                topic_elem = copy.deepcopy(templates['topic'])
                slide.shapes._spTree.append(topic_elem)
                shape_elements.append(slide.shapes[-1])
            
            if templates.get('subtopic') is not None and 'subtopic' in data:
                subtopic_elem = copy.deepcopy(templates['subtopic'])
                slide.shapes._spTree.append(subtopic_elem)
                shape_elements.append(slide.shapes[-1])
            
            if templates.get('text') is not None:
                text_elem = copy.deepcopy(templates['text'])
                # If no subtopic, move the text box to where the subtopic would have been
                if 'subtopic' not in data and templates.get('subtopic') is not None:
                    sub_xml = templates['subtopic']
                    # Search for offset element in subtopic template
                    # The namespace for 'a' is usually needed
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    sub_off = sub_xml.find('.//a:off', namespaces=ns)
                    text_off = text_elem.find('.//a:off', namespaces=ns)
                
                    if sub_off is not None and text_off is not None:
                        # Move text to subtopic's Y coordinate
                        text_off.set('y', sub_off.get('y'))
            
                slide.shapes._spTree.append(text_elem)
                shape_elements.append(slide.shapes[-1])
        
            # Insert text and subtexts
            for shape in shape_elements:
                if shape.has_text_frame:
                    shape_text = shape.text
                    if 'Text goes here' in shape_text:
                        if 'text' in data:
                            replace_text_preserve_format(shape, data['text'], font_color=RGBColor(255, 255, 255))
            
                # Check for topic/subtopic either as standalone shape or inside group
                is_group = (getattr(shape, 'shape_type', None) == 6)
                targets = shape.shapes if is_group else [shape]
            
                for target in targets:
                    if getattr(target, 'has_text_frame', False):
                        txt = target.text.strip().lower()
                        if 'topic' in data or 'subtopic' in data:
                            # print(f"DEBUG: Comparing txt='{txt}' with data keys {list(data.keys())}")
                            pass
                        if txt == 'topic' and 'topic' in data:
                            replace_text_preserve_format(target, data['topic'], center=True)
                        elif txt == 'subtopic' and 'subtopic' in data:
                            replace_text_preserve_format(target, data['subtopic'], center=True)
                            target.text_frame.word_wrap = False
                            # Resize logic: if it's a group, resize all children. If standalone, resize it.
                            estimated_width = len(data['subtopic']) * 300000 + 400000
                            if estimated_width > shape.width:
                                diff = estimated_width - shape.width
                                shape.width = estimated_width
                                if is_group:
                                    for child in shape.shapes:
                                        child.width = child.width + diff
                                else:
                                    pass # Shape width already updated
                            
            # Handle picture placeholder(s)
            images = section.get('images', [])
            pic_ph = None
            for shape in slide.shapes:
                if getattr(shape, 'is_placeholder', False) and shape.placeholder_format.type == 18:
                    pic_ph = shape
                    break
        
            if pic_ph is not None and images:
                # Capture properties of the placeholder before it's potentially replaced/removed
                left, top, width, height = pic_ph.left, pic_ph.top, pic_ph.width, pic_ph.height
                num_images = len(images)
            
                if num_images > 1:
                    # Remove the placeholder to avoid overlap with newly added pictures
                    pic_ph.element.getparent().remove(pic_ph.element)
                
                    gap = Inches(0.1)
                    total_gap = gap * (num_images - 1)
                    img_width = (width - total_gap) / num_images
                
                    for i, img_blob in enumerate(images):
                        # Load image to get original dimensions for aspect ratio
                        with Image.open(io.BytesIO(img_blob)) as pil_img:
                            orig_w, orig_h = pil_img.size
                    
                        slot_w = img_width
                        slot_h = height
                    
                        # Calculate fit dimensions
                        slot_aspect = slot_w / slot_h
                        img_aspect = orig_w / orig_h
                    
                        if img_aspect > slot_aspect:
                            # Image is wider than slot relative to height
                            draw_w = slot_w
                            draw_h = slot_w / img_aspect
                        else:
                            # Image is taller than slot relative to width
                            draw_h = slot_h
                            draw_w = slot_h * img_aspect
                    
                        # Center within slot
                        final_left = int(left + (i * (slot_w + gap)) + (slot_w - draw_w) / 2)
                        final_top = int(top + (slot_h - draw_h) / 2)
                    
                        slide.shapes.add_picture(
                            io.BytesIO(img_blob), 
                            final_left, 
                            final_top, 
                            width=int(draw_w),
                            height=int(draw_h)
                        )
                else:
                    # Single image: use standard placeholder insertion (handles aspect ratio better)
                    pic_ph.insert_picture(io.BytesIO(images[0]))
            elif pic_ph is not None:
                # Hide/Remove placeholder if no images are present
                pic_ph.element.getparent().remove(pic_ph.element)
                    
            print(f"Updated {layout.name} with text and {len(images)} images.")

        elif layout.name in ('LAYOUT_sst_activity_page_01', 'LAYOUT_sst_activity_page_02'):
            templates = sst_content_templates.get(layout.name, {})
        
            # Inject static elements
            for static_elem in templates.get('static_elements', []):
                elem_copy = remove_locks(copy.deepcopy(static_elem))
                copy_image_rels(elem_copy, layout.part, slide.part)
                slide.shapes._spTree.append(elem_copy)

            # Handle text injection and dynamic resizing
            text_content = []
            for entry in section['content']:
                content_obj = entry[0] if isinstance(entry, tuple) else entry
                text_content.append(content_obj)

            if templates.get('text') is not None:
                text_elem = copy.deepcopy(templates['text'])
                slide.shapes._spTree.append(text_elem)
                shape = slide.shapes[-1]
                # Force black text for activity boxes if they are yellow
                replace_text_preserve_format(shape, text_content, font_color=RGBColor(0, 0, 0))
            
                # Estimate height needed for 36pt font
                chars_per_line = max(20, int((shape.width / 914400) * 3.5))
            
                line_count = 0
                for line_item in text_content:
                    # Get plain text length for estimation
                    if isinstance(line_item, list):
                        line_str = "".join(p['value'] if p['type'] == 'text' else '    ' for p in line_item)
                    else:
                        line_str = str(line_item)
                    line_count += (len(line_str) // chars_per_line) + 1
                
                if line_count > 1:
                    # 36pt font is ~0.5 inch high. With spacing, we need ~0.7 inch per line.
                    # EMU: 1 inch = 914400. 0.7 inch = ~640,000 EMUs.
                    extra_h = (line_count - 1) * 650000 
                    old_h = shape.height
                    shape.height += extra_h
                    new_h = shape.height
                
                    # Maintain corner roundness
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    gd = shape.element.find('.//a:gd[@name="adj"]', namespaces=ns)
                    if gd is not None:
                        try:
                            fmla = gd.get('fmla')
                            if 'val' in fmla:
                                old_adj = int(fmla.split()[-1])
                                new_adj = int(old_adj * (old_h / new_h))
                                gd.set('fmla', f'val {new_adj}')
                        except:
                            pass
                
                    # Move everything below it down (e.g. discussion group)
                    for other_shape in slide.shapes:
                        if other_shape != shape and other_shape.top > shape.top:
                             # Use a bit more buffer when moving things down
                             other_shape.top += extra_h + 100000

            # Handle picture placeholder
            images = section.get('images', [])
            pic_ph = None
            for shp in slide.shapes:
                if getattr(shp, 'is_placeholder', False) and shp.placeholder_format.type == 18:
                    pic_ph = shp
                    break
        
            if pic_ph is not None and images:
                # Capture properties of the placeholder
                left, top, total_width, total_height = pic_ph.left, pic_ph.top, pic_ph.width, pic_ph.height
                num_images = len(images)
            
                if num_images > 1:
                    # Remove the placeholder to avoid overlap
                    pic_ph.element.getparent().remove(pic_ph.element)
                
                    # We want square slots. The width of each slot is total_width / num_images.
                    # But we should also consider the height. 
                    # Let's use the smaller of (total_width/num_images) and total_height as the square side.
                    gap = Inches(0.1)
                    slot_side = min((total_width - (gap * (num_images-1))) / num_images, total_height)
                
                    for i, img_blob in enumerate(images):
                        with Image.open(io.BytesIO(img_blob)) as pil_img:
                            orig_w, orig_h = pil_img.size
                    
                        img_aspect = orig_w / orig_h
                    
                        # Calculate dimensions to fit inside the square slot_side x slot_side
                        if img_aspect > 1: # Wide image
                            draw_w = slot_side
                            draw_h = slot_side / img_aspect
                        else: # Tall or square image
                            draw_h = slot_side
                            draw_w = slot_side * img_aspect
                        
                        # Center within the slot
                        slot_left = left + i * (slot_side + gap)
                        # For layout 01, we center vertically in the placeholder's original vertical space
                        final_left = int(slot_left + (slot_side - draw_w) / 2)
                        final_top = int(top + (total_height - draw_h) / 2)
                    
                        slide.shapes.add_picture(
                            io.BytesIO(img_blob), 
                            final_left, 
                            final_top, 
                            width=int(draw_w),
                            height=int(draw_h)
                        )
                else:
                    # Single image: use standard placeholder insertion
                    slide.shapes._spTree.append(pic_ph.element)
                    pic_ph.insert_picture(io.BytesIO(images[0]))
            elif pic_ph is not None:
                pic_ph.element.getparent().remove(pic_ph.element)

            print(f"Updated {layout.name} with activity content and {len(images)} images.")

        elif layout.name == 'LAYOUT_final_quiz_page':
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    shape.text = "Final Quiz!"
                    break
            print("Inserted LAYOUT_final_quiz_page (Welcome Slide)")

            qa_pairs = []
            for entry in section['content']:
                line = get_text(entry)
                if ":" in line:
                    parts = line.split(":", 1)
                    q_text = parts[0].strip()
                    a_text = parts[1].strip()
                    if q_text.lower() == 'question' and a_text.lower() == 'answer':
                        continue
                    qa_pairs.append((q_text, a_text))
        
            layout_q = get_layout('final_quiz_page_q')
            layout_a = get_layout('final_quiz_page_a')
        
            for q, a in qa_pairs:
                if layout_q:
                    s_q = prs.slides.add_slide(layout_q)
                    for shape in s_q.shapes:
                        if shape.is_placeholder and shape.placeholder_format.idx == 0:
                            shape.text = q
                if layout_a:
                    s_a = prs.slides.add_slide(layout_a)
                    for shape in s_a.shapes:
                        if shape.is_placeholder and shape.placeholder_format.idx == 0:
                            shape.text = a
            print(f"Generated {len(qa_pairs)} Question and Answer slide pairs.")
        elif layout.name == 'LAYOUT_sst_notedown_page':
            templates = sst_content_templates.get(layout.name, {})
            if templates.get('text') is not None:
                text_elem = copy.deepcopy(templates['text'])
                slide.shapes._spTree.append(text_elem)
                shape = slide.shapes[-1]
            
                processed_content = []
                for entry in section['content']:
                    line = get_text(entry)
                    clean_line = line.strip()
                    if clean_line.startswith(("•", "-", "*")):
                        clean_line = "❖ " + clean_line.lstrip("•-*").strip()
                    processed_content.append(clean_line)
            
                replace_text_preserve_format(shape, processed_content)
                print(f"Updated {layout.name} Body text with preserved formatting and spacing.")

        elif layout.name == 'LAYOUT_sst_discussion_page':
            templates = sst_content_templates.get(layout.name, {})
            # Inject all static elements (groups, pictures, title) from layout
            for static_elem in templates.get('static_elements', []):
                elem_copy = remove_locks(copy.deepcopy(static_elem))
                copy_image_rels(elem_copy, layout.part, slide.part)
                slide.shapes._spTree.append(elem_copy)
            if templates.get('question1') is not None:
                # Parse all content lines from the section
                q1_lines = []
                for entry in section['content']:
                    line = get_text(entry)
                    if ":" in line:
                        parts = line.split(":", 1)
                        if parts[0].strip().lower() == 'question1':
                            q1_lines.append(parts[1].strip())
                        else:
                            q1_lines.append(line.strip())
                    elif line.strip():
                        q1_lines.append(line.strip())
            
                q1_elem = copy.deepcopy(templates['question1'])
                slide.shapes._spTree.append(q1_elem)
                shape = slide.shapes[-1]
                if q1_lines:
                    replace_text_preserve_format(shape, q1_lines, font_color=RGBColor(255, 255, 255))
                    # If single line, remove bullet formatting from paragraphs
                    if len(q1_lines) == 1:
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        for p_elem in shape.text_frame._txBody.findall('.//a:p', namespaces=ns):
                            pPr = p_elem.find('a:pPr', namespaces=ns)
                            if pPr is not None:
                                for tag in ('a:buFont', 'a:buChar', 'a:buAutoNum'):
                                    el = pPr.find(tag, namespaces=ns)
                                    if el is not None:
                                        pPr.remove(el)
                                pPr.set('marL', '0')
                                pPr.set('indent', '0')
                    print(f"Updated {layout.name} question1 -> {len(q1_lines)} lines")
                else:
                    replace_text_preserve_format(shape, "question1")  # Keep original if missing
                    print(f"Inserted {layout.name} with default text")

        elif layout.name == 'LAYOUT_sst_homework_page':
            templates = sst_content_templates.get(layout.name, {})
            # Inject all static elements (groups, text boxes, etc) from layout
            for static_elem in templates.get('static_elements', []):
                elem_copy = remove_locks(copy.deepcopy(static_elem))
                copy_image_rels(elem_copy, layout.part, slide.part)
                slide.shapes._spTree.append(elem_copy)
            print(f"Inserted {layout.name} with all background elements")

        # After processing the section and adding elements, check if we need to overlay SYR
        if section.get('has_syr'):
            layout_syr = get_layout('syr')
            if layout_syr:
                syr_templates = sst_content_templates.get('LAYOUT_syr', {})
                for static_elem in syr_templates.get('static_elements', []):
                    elem_copy = remove_locks(copy.deepcopy(static_elem))
                    copy_image_rels(elem_copy, layout_syr.part, slide.part)
                    slide.shapes._spTree.append(elem_copy)
                print(f"Overlaid LAYOUT_syr on current slide.")

        # Check if we need to overlay Ask Question
        ask_q_text = section.get('ask_question_text')
        if ask_q_text:
            layout_ask = get_layout('ask_question')
            if layout_ask:
                ask_templates = sst_content_templates.get('LAYOUT_ask_question', {})
                for static_elem in ask_templates.get('static_elements', []):
                    elem_copy = remove_locks(copy.deepcopy(static_elem))
                    copy_image_rels(elem_copy, layout_ask.part, slide.part)
                    slide.shapes._spTree.append(elem_copy)
                
                # After appending all Ask Question elements, find the group and update its text
                for shp in slide.shapes:
                    if shp.shape_type == 6:  # GROUP
                        for sub in shp.shapes:
                            if getattr(sub, 'has_text_frame', False):
                                if 'write question here' in sub.text.lower():
                                    replace_text_preserve_format(sub, ask_q_text)
                                    # Add 10% padding to left and right inside the text frame
                                    margin = int(sub.width * 0.10)
                                    sub.text_frame.margin_left = margin
                                    sub.text_frame.margin_right = margin
                                    # Center if short (one line), left align if long (multi-line)
                                    align = PP_ALIGN.CENTER if len(ask_q_text) <= 45 else PP_ALIGN.LEFT
                                    for p in sub.text_frame.paragraphs:
                                        p.alignment = align
                print(f"Overlaid LAYOUT_ask_question with text: '{ask_q_text}'")

    prs.save(output_path)

if __name__ == "__main__":
    generate_ppt("content.docx", "template.pptx", "Generated_Presentation.pptx")
    print("DONE -- PPT Generated")
