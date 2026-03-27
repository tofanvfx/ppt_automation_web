"""Check buChar placement in each paragraph XML."""
from pptx import Presentation
from lxml import etree

prs = Presentation('test_run/sub_bullet_test.pptx')
sl = prs.slides[0]
sh = [s for s in sl.shapes if s.has_text_frame and 'Google Shape' in s.name][0]
tf = sh.text_frame
a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

for p in tf.paragraphs:
    if not p.text.strip():
        continue
    # Full pPr XML
    pPr = p._p.find(f'{{{a_ns}}}pPr')
    pPr_xml = etree.tostring(pPr, pretty_print=True).decode() if pPr is not None else '(no pPr)'
    print(f"Level={p.level} text={p.text[:30]:30s}: {pPr_xml[:200]}")
    print("---")
