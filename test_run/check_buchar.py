from pptx import Presentation

prs = Presentation('test_run/sub_bullet_test.pptx')
sl = prs.slides[0]
sh = [s for s in sl.shapes if s.has_text_frame and 'Google Shape' in s.name][0]
tf = sh.text_frame
for p in tf.paragraphs:
    if not p.text.strip():
        continue
    bu_chars = [c.get('char') for c in p._p.iter() if c.tag.endswith('}buChar')]
    bu_nones = [c.tag for c in p._p.iter() if c.tag.endswith('}buNone')]
    print(f"Level={p.level} text={p.text[:35]!r} buChar={bu_chars} buNone={bu_nones}")
